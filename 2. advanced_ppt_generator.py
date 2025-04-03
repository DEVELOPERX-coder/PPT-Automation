import os
import re
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL

class AdvancedPPTGenerator:
    def __init__(self, template_path=None):
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        # Standard slide layouts dictionary
        self.slide_layouts = {
            'title': 0,              # Title Slide
            'title_content': 1,      # Title and Content
            'section': 2,            # Section Header
            'two_content': 3,        # Two Content
            'comparison': 4,         # Comparison
            'title_only': 5,         # Title Only
            'blank': 6               # Blank
        }
        
        # Keep track of the current slide for animations
        self.current_slide = None
        self.slide_index = 0
    
    def set_slide_size(self, width, height):
        """Set slide size in inches (widescreen 16:9 = 13.33 x 7.5)"""
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)
    
    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        # Remove the # if present
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def add_slide(self, layout_type='blank'):
        """Add a new slide with the specified layout"""
        if layout_type in self.slide_layouts:
            layout_idx = self.slide_layouts[layout_type]
        else:
            print(f"Warning: Layout '{layout_type}' not found. Using 'blank' instead.")
            layout_idx = self.slide_layouts['blank']
            
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)
        self.current_slide = slide
        self.slide_index += 1
        return slide
    
    def set_background_color(self, slide, color_hex):
        """Set slide background to a solid color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        rgb = self.hex_to_rgb(color_hex)
        fill.fore_color.rgb = RGBColor(*rgb)
    
    def add_textbox(self, slide, text, left, top, width, height, 
               font_name=None, font_size=None, bold=False, italic=False,
               color_hex=None, alignment=PP_ALIGN.LEFT):
        """Add a text box with formatting"""
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        
        if font_name or font_size or bold or italic or color_hex:
            # Check if there are any runs before trying to access them
            if not p.runs:
                run = p.add_run()
                run.text = text
            else:
                run = p.runs[0]
                
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            if color_hex:
                rgb = self.hex_to_rgb(color_hex)
                run.font.color.rgb = RGBColor(*rgb)
                
        return textbox, p
    
    def add_shape(self, slide, shape_type, left, top, width, height, 
                 fill_color_hex=None, line_color_hex=None, line_width=None):
        """Add a shape to the slide"""
        shape = slide.shapes.add_shape(
            shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        
        if fill_color_hex:
            rgb = self.hex_to_rgb(fill_color_hex)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb)
            
        if line_color_hex:
            rgb = self.hex_to_rgb(line_color_hex)
            shape.line.color.rgb = RGBColor(*rgb)
            
        if line_width:
            shape.line.width = Pt(line_width)
            
        return shape
    
    def add_rectangle(self, slide, left, top, width, height, 
                     fill_color_hex=None, line_color_hex=None, line_width=None):
        """Add a rectangle to the slide"""
        return self.add_shape(
            slide, MSO_SHAPE.RECTANGLE, 
            left, top, width, height,
            fill_color_hex, line_color_hex, line_width
        )
    
    def add_code_block(self, slide, code, left, top, width, height, 
                      language='cpp', font_name='Consolas', font_size=11,
                      text_color_hex='#03DAC6', bg_color_hex='#252525'):
        """Add a code block with syntax highlighting simulation"""
        # Create the rectangle background
        shape = self.add_rectangle(
            slide, left, top, width, height,
            fill_color_hex=bg_color_hex
        )
        
        # Add the code as a textbox
        textbox, p = self.add_textbox(
            slide, "", left + 0.1, top + 0.1, width - 0.2, height - 0.2,
            font_name=font_name, font_size=font_size
        )
        
        # Basic syntax highlighting for C++
        # Note: This is a simplified version; a full implementation would use regex patterns
        code_lines = code.split('\n')
        
        tf = textbox.text_frame
        tf.clear()  # Clear default paragraph
        
        for line in code_lines:
            para = tf.add_paragraph()
            
            # This is a simplified approach - in a real application,
            # you would use regex patterns to properly tokenize the code
            parts = re.split(r'(\b(?:int|char|for|if|else|return|void|double|float|std|string)\b|".*?"|//.*$)', line)
            
            for part in parts:
                if part:
                    run = para.add_run()
                    run.text = part
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                    
                    # Color keywords
                    if part.strip() in ['int', 'char', 'for', 'if', 'else', 'return', 'void', 'double', 'float']:
                        rgb = self.hex_to_rgb('#FF79C6')  # Purple for keywords
                    elif part.strip() in ['std', 'string']:
                        rgb = self.hex_to_rgb('#8BE9FD')  # Cyan for library names
                    elif part.startswith('"') and part.endswith('"'):
                        rgb = self.hex_to_rgb('#F1FA8C')  # Yellow for strings
                    elif part.startswith('//'):
                        rgb = self.hex_to_rgb('#6272A4')  # Grayish for comments
                    else:
                        rgb = self.hex_to_rgb(text_color_hex)  # Default text color
                        
                    run.font.color.rgb = RGBColor(*rgb)
        
        return textbox
    
    def add_table(self, slide, rows, cols, left, top, width, height,
                 header_fill_hex=None, alt_row_fill_hex=None, 
                 header_text_hex='#FFFFFF'):
        """Add a table to the slide"""
        table = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        ).table
        
        # Set header row formatting
        if header_fill_hex:
            header_rgb = self.hex_to_rgb(header_fill_hex)
            for cell in table.rows[0].cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*header_rgb)
                
                # Set text color for header
                header_text_rgb = self.hex_to_rgb(header_text_hex)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(*header_text_rgb)
        
        # Set alternating row colors
        if alt_row_fill_hex:
            alt_rgb = self.hex_to_rgb(alt_row_fill_hex)
            for i in range(2, rows, 2):  # Start from row 2 (0-indexed), every other row
                for cell in table.rows[i].cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*alt_rgb)
        
        return table
    
    def set_table_cell_text(self, table, row, col, text, 
                      font_name=None, font_size=None, bold=False,
                      color_hex=None, alignment=PP_ALIGN.LEFT):
        """Set text and formatting for a table cell"""
        cell = table.cell(row, col)
        cell.text = text
        
        # Format the text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = alignment
            
            # Check if there are any runs before trying to access them
            if not paragraph.runs:
                run = paragraph.add_run()
                run.text = text
            else:
                if not paragraph.runs:
                    run = paragraph.add_run()
                    run.text = text_content  # Use appropriate text here
                else:
                    run = paragraph.runs[0]
                
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = Pt(font_size)
            run.font.bold = bold
            if color_hex:
                rgb = self.hex_to_rgb(color_hex)
                run.font.color.rgb = RGBColor(*rgb)
    
    def add_animation_note(self, slide, animation_type, element_desc, 
                         timing=None, direction=None, duration=None, delay=None):
        """Add animation instructions as presenter notes"""
        # This doesn't actually create animations, but adds notes that a person can follow
        # to manually add animations in PowerPoint
        if not hasattr(slide, 'animation_notes'):
            slide.animation_notes = []
            
        note = f"Animation: {animation_type} for '{element_desc}'"
        if timing:
            note += f", Timing: {timing}"
        if direction:
            note += f", Direction: {direction}"
        if duration:
            note += f", Duration: {duration}s"
        if delay:
            note += f", Delay: {delay}s"
            
        slide.animation_notes.append(note)
        
        # Also add to slide notes
        if not slide.has_notes_slide:
            slide.notes_slide
            
        if slide.notes_slide.notes_text_frame.text:
            slide.notes_slide.notes_text_frame.text += f"\n{note}"
        else:
            slide.notes_slide.notes_text_frame.text = note
    
    def apply_gradient_fill_note(self, slide, element_desc, color1, color2, direction="horizontal"):
        """Add notes about applying a gradient fill to an element"""
        note = f"Apply gradient fill to '{element_desc}': From {color1} to {color2}, Direction: {direction}"
        
        # Add to slide notes
        if not slide.has_notes_slide:
            slide.notes_slide
            
        if slide.notes_slide.notes_text_frame.text:
            slide.notes_slide.notes_text_frame.text += f"\n{note}"
        else:
            slide.notes_slide.notes_text_frame.text = note
    
    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        print(f"Presentation saved as {filename}")
        
        # Print any animation notes that need to be manually implemented
        print("\nManual Animation/Formatting Steps Required:")
        for i, slide in enumerate(self.prs.slides):
            if hasattr(slide, 'animation_notes'):
                print(f"\nSlide {i+1}:")
                for note in slide.animation_notes:
                    print(f"  - {note}")


class PPTSpecParser:
    """Parse PowerPoint specification files into automation commands"""
    
    def __init__(self, spec_file=None):
        self.spec_data = None
        if spec_file:
            self.load_spec(spec_file)
    
    def load_spec(self, spec_file):
        """Load specification from YAML file"""
        with open(spec_file, 'r') as f:
            self.spec_data = yaml.safe_load(f)
        return self.spec_data
    
    def generate_presentation(self, output_file):
        """Generate presentation based on loaded specification"""
        if not self.spec_data:
            raise ValueError("No specification data loaded")
            
        # Create presentation generator
        ppt = AdvancedPPTGenerator()
        
        # Set slide size
        if 'slide_size' in self.spec_data:
            size = self.spec_data['slide_size']
            ppt.set_slide_size(size['width'], size['height'])
            
        # Process slides
        for slide_spec in self.spec_data['slides']:
            self._process_slide(ppt, slide_spec)
            
        # Save presentation
        ppt.save(output_file)
        
    def _process_slide(self, ppt, slide_spec):
        """Process a single slide specification"""
        # Create slide
        layout = slide_spec.get('layout', 'blank')
        slide = ppt.add_slide(layout)
        
        # Set background
        if 'background' in slide_spec:
            bg = slide_spec['background']
            if 'color' in bg:
                ppt.set_background_color(slide, bg['color'])
        
        # Add title if present
        if 'title' in slide_spec:
            title = slide_spec['title']
            left = title.get('left', 1)
            top = title.get('top', 0.8)
            width = title.get('width', 10)
            height = title.get('height', 1)
            
            ppt.add_textbox(
                slide, title['text'], left, top, width, height,
                font_name=title.get('font_name', 'Segoe UI Light'),
                font_size=title.get('font_size', 40),
                bold=title.get('bold', True),
                color_hex=title.get('color', '#00FFFF'),
                alignment=PP_ALIGN.LEFT
            )
        
        # Process elements
        if 'elements' in slide_spec:
            for element in slide_spec['elements']:
                self._process_element(ppt, slide, element)
        
        # Process animations
        if 'animations' in slide_spec:
            for animation in slide_spec['animations']:
                ppt.add_animation_note(
                    slide,
                    animation['type'],
                    animation['element'],
                    animation.get('timing'),
                    animation.get('direction'),
                    animation.get('duration'),
                    animation.get('delay')
                )
    
    def _process_element(self, ppt, slide, element):
        """Process a single slide element"""
        element_type = element['type']
        
        if element_type == 'textbox':
            left = element.get('left', 1)
            top = element.get('top', 1)
            width = element.get('width', 10)
            height = element.get('height', 1)
            
            ppt.add_textbox(
                slide, element['text'], left, top, width, height,
                font_name=element.get('font_name'),
                font_size=element.get('font_size'),
                bold=element.get('bold', False),
                italic=element.get('italic', False),
                color_hex=element.get('color'),
                alignment=PP_ALIGN.LEFT if 'alignment' not in element else PP_ALIGN.CENTER 
                if element['alignment'] == 'center' else PP_ALIGN.RIGHT
            )
            
            # Note about gradient if specified
            if 'gradient' in element:
                g = element['gradient']
                ppt.apply_gradient_fill_note(
                    slide, element['text'],
                    g['color1'], g['color2'],
                    g.get('direction', 'horizontal')
                )
                
        elif element_type == 'shape':
            shape_type = element.get('shape_type', 'rectangle')
            left = element.get('left', 1)
            top = element.get('top', 1)
            width = element.get('width', 3)
            height = element.get('height', 1)
            
            if shape_type == 'rectangle':
                ppt.add_rectangle(
                    slide, left, top, width, height,
                    fill_color_hex=element.get('fill_color'),
                    line_color_hex=element.get('line_color'),
                    line_width=element.get('line_width')
                )
                
        elif element_type == 'code':
            left = element.get('left', 1)
            top = element.get('top', 3)
            width = element.get('width', 10)
            height = element.get('height', 3)
            
            ppt.add_code_block(
                slide, element['code'], left, top, width, height,
                language=element.get('language', 'cpp'),
                font_name=element.get('font_name', 'Consolas'),
                font_size=element.get('font_size', 14),
                text_color_hex=element.get('text_color', '#03DAC6'),
                bg_color_hex=element.get('bg_color', '#252525')
            )
            
        elif element_type == 'table':
            left = element.get('left', 1)
            top = element.get('top', 3)
            width = element.get('width', 10)
            height = element.get('height', 4)
            rows = element.get('rows', 3)
            cols = element.get('cols', 3)
            
            table = ppt.add_table(
                slide, rows, cols, left, top, width, height,
                header_fill_hex=element.get('header_fill'),
                alt_row_fill_hex=element.get('alt_row_fill')
            )
            
            # Add cell content if provided
            if 'cells' in element:
                for cell in element['cells']:
                    row = cell['row']
                    col = cell['col']
                    ppt.set_table_cell_text(
                        table, row, col, cell['text'],
                        font_name=cell.get('font_name'),
                        font_size=cell.get('font_size'),
                        bold=cell.get('bold', False),
                        color_hex=cell.get('color')
                    )


def convert_text_spec_to_yaml(text_spec, output_file):
    """
    Convert a text-based specification (like your PDF) into a YAML format
    that our parser can understand. This is a simplified example.
    """
    # This would be a complex parser for your specific format
    # For demonstration, I'll just create a basic YAML structure
    
    yaml_spec = {
        'slide_size': {'width': 13.33, 'height': 7.5},  # 16:9 ratio
        'slides': []
    }
    
    # Very basic parsing - in practice you'd need a much more robust parser
    current_slide = None
    lines = text_spec.split('\n')
    
    for line in lines:
        if line.strip().startswith('SLIDE '):
            # New slide found
            if current_slide:
                yaml_spec['slides'].append(current_slide)
                
            slide_title = line.split(':', 1)[1].strip() if ':' in line else ""
            current_slide = {
                'layout': 'blank',
                'background': {'color': '#1A1A1A'},  # Dark charcoal
                'title': {
                    'text': slide_title,
                    'left': 1,
                    'top': 0.8,
                    'width': 10,
                    'height': 1,
                    'font_name': 'Segoe UI Light',
                    'font_size': 40,
                    'bold': True,
                    'color': '#00FFFF'  # Cyan
                },
                'elements': [],
                'animations': []
            }
    
    # Add the last slide if there is one
    if current_slide:
        yaml_spec['slides'].append(current_slide)
    
    # Write to YAML file
    with open(output_file, 'w') as f:
        yaml.dump(yaml_spec, f, default_flow_style=False)
    
    return yaml_spec


def create_demo_spec():
    """Create a demonstration specification file for the C++ Strings presentation"""
    spec = {
        'slide_size': {'width': 13.33, 'height': 7.5},  # 16:9 ratio
        'slides': [
            # Title slide
            {
                'layout': 'blank',
                'background': {'color': '#1A1A1A'},  # Dark charcoal
                'title': {
                    'text': '🧵 STRINGS IN C++',
                    'left': 3.5,
                    'top': 3,
                    'width': 7,
                    'height': 1.5,
                    'font_name': 'Segoe UI Light',
                    'font_size': 48,
                    'bold': True,
                    'color': '#00FFFF',  # Cyan
                    'alignment': 'center'
                },
                'elements': [
                    {
                        'type': 'textbox',
                        'text': 'Arrays With Character(s)',
                        'left': 3.5,
                        'top': 4.5,
                        'width': 6,
                        'height': 1,
                        'font_name': 'Segoe UI',
                        'font_size': 32,
                        'color': '#BB86FC',  # Light Purple
                        'alignment': 'center'
                    },
                    {
                        'type': 'textbox',
                        'text': '#DSAin45 - Day 4',
                        'left': 5.5,
                        'top': 5.5,
                        'width': 3,
                        'height': 0.6,
                        'font_name': 'Segoe UI',
                        'font_size': 18,
                        'color': '#E0E0E0',  # Light Gray
                        'alignment': 'center'
                    }
                ],
                'animations': [
                    {
                        'type': 'Float In',
                        'element': 'Title',
                        'direction': 'From Left',
                        'duration': 1.5
                    },
                    {
                        'type': 'Fade',
                        'element': 'Subtitle',
                        'delay': 0.3,
                        'duration': 1
                    }
                ]
            },
            
            # What are Strings slide
            {
                'layout': 'blank',
                'background': {'color': '#1A1A1A'},
                'title': {
                    'text': 'What Exactly ARE Strings?',
                    'left': 1,
                    'top': 0.8,
                    'width': 10,
                    'height': 1,
                    'font_name': 'Segoe UI Light',
                    'font_size': 40,
                    'bold': True,
                    'gradient': {
                        'color1': '#00FFFF',  # Cyan
                        'color2': '#00BFFF',  # Light Blue
                        'direction': 'horizontal'
                    }
                },
                'elements': [
                    {
                        'type': 'textbox',
                        'text': 'At their core, strings are sequences of characters. But in C++, there are two main ways to represent them:',
                        'left': 1,
                        'top': 2,
                        'width': 10,
                        'height': 0.8,
                        'font_name': 'Segoe UI',
                        'font_size': 20,
                        'color': '#E0E0E0'  # Light Gray
                    },
                    {
                        'type': 'shape',
                        'shape_type': 'rectangle',
                        'left': 1,
                        'top': 3,
                        'width': 5.5,
                        'height': 5,
                        'fill_color': '#2D2D2D',  # Dark Gray
                        'line_color': '#00FFFF',  # Cyan
                        'line_width': 1
                    },
                    {
                        'type': 'textbox',
                        'text': '1. C-style Strings (char arrays)',
                        'left': 1.25,
                        'top': 3.2,
                        'width': 5,
                        'height': 0.7,
                        'font_name': 'Segoe UI',
                        'font_size': 22,
                        'bold': True,
                        'color': '#00FFFF'  # Cyan
                    },
                    {
                        'type': 'code',
                        'code': 'char greeting[] = "Hello"; // Compiler adds null terminator \'\\0\'',
                        'left': 1.25,
                        'top': 4,
                        'width': 5,
                        'height': 0.6,
                        'font_name': 'Consolas',
                        'font_size': 16,
                        'text_color': '#03DAC6',  # Light Green
                        'bg_color': '#252525'  # Very Dark Gray
                    },
                    {
                        'type': 'shape',
                        'shape_type': 'rectangle',
                        'left': 7.5,
                        'top': 3,
                        'width': 5.5,
                        'height': 5,
                        'fill_color': '#2D2D2D',  # Dark Gray
                        'line_color': '#BB86FC',  # Light Purple
                        'line_width': 1
                    },
                    {
                        'type': 'textbox',
                        'text': '2. C++ std::string',
                        'left': 7.75,
                        'top': 3.2,
                        'width': 5,
                        'height': 0.7,
                        'font_name': 'Segoe UI',
                        'font_size': 22,
                        'bold': True,
                        'color': '#BB86FC'  # Light Purple
                    },
                    {
                        'type': 'code',
                        'code': 'std::string greeting = "Hello";',
                        'left': 7.75,
                        'top': 4,
                        'width': 5,
                        'height': 0.6,
                        'font_name': 'Consolas',
                        'font_size': 16,
                        'text_color': '#03DAC6',  # Light Green
                        'bg_color': '#252525'  # Very Dark Gray
                    }
                ],
                'animations': [
                    {
                        'type': 'Split',
                        'element': 'Title',
                        'direction': 'Horizontal Out',
                        'duration': 1
                    },
                    {
                        'type': 'Fade',
                        'element': 'Definition paragraph',
                        'duration': 0.7
                    },
                    {
                        'type': 'Wipe',
                        'element': 'C-style section',
                        'direction': 'From Left',
                        'duration': 0.8
                    },
                    {
                        'type': 'Wipe',
                        'element': 'std::string section',
                        'direction': 'From Right',
                        'duration': 0.8
                    }
                ]
            },
            
            # Memory Layout slide
            {
                'layout': 'blank',
                'background': {'color': '#1A1A1A'},
                'title': {
                    'text': 'String Memory Layout',
                    'left': 1,
                    'top': 0.8,
                    'width': 10,
                    'height': 1,
                    'font_name': 'Segoe UI Light',
                    'font_size': 40,
                    'bold': True,
                    'gradient': {
                        'color1': '#00FFFF',  # Cyan
                        'color2': '#00BFFF',  # Light Blue
                        'direction': 'horizontal'
                    }
                },
                'elements': [
                    {
                        'type': 'textbox',
                        'text': 'C-style String Memory:',
                        'left': 1,
                        'top': 2,
                        'width': 5,
                        'height': 0.7,
                        'font_name': 'Segoe UI',
                        'font_size': 24,
                        'bold': True,
                        'color': '#00FFFF'  # Cyan
                    },
                    {
                        'type': 'textbox',
                        'text': 'std::string Memory:',
                        'left': 7,
                        'top': 2,
                        'width': 5,
                        'height': 0.7,
                        'font_name': 'Segoe UI',
                        'font_size': 24,
                        'bold': True,
                        'color': '#BB86FC'  # Light Purple
                    }
                ],
                'animations': [
                    {
                        'type': 'Fade',
                        'element': 'Title',
                        'duration': 0.7
                    },
                    {
                        'type': 'Fade',
                        'element': 'C-style heading',
                        'duration': 0.5
                    },
                    {
                        'type': 'Fade',
                        'element': 'std::string heading',
                        'duration': 0.5
                    }
                ]
            },
            
            # Operations & Complexity Table
            {
                'layout': 'blank',
                'background': {'color': '#1A1A1A'},
                'title': {
                    'text': 'String Operations & Complexity',
                    'left': 1,
                    'top': 0.8,
                    'width': 10,
                    'height': 1,
                    'font_name': 'Segoe UI Light',
                    'font_size': 40,
                    'bold': True,
                    'gradient': {
                        'color1': '#00FFFF',  # Cyan
                        'color2': '#00BFFF',  # Light Blue
                        'direction': 'horizontal'
                    }
                },
                'elements': [
                    {
                        'type': 'table',
                        'left': 0.5,
                        'top': 2,
                        'width': 12,
                        'height': 4.5,
                        'rows': 8,
                        'cols': 3,
                        'header_fill': '#3D3D3D',  # Dark Gray
                        'alt_row_fill': '#333333',  # Slightly lighter gray for alternating rows
                        'cells': [
                            # Headers
                            {'row': 0, 'col': 0, 'text': 'Operation', 'font_name': 'Segoe UI', 'font_size': 20, 'bold': True, 'color': '#FFFFFF'},
                            {'row': 0, 'col': 1, 'text': 'C-style', 'font_name': 'Segoe UI', 'font_size': 20, 'bold': True, 'color': '#00FFFF'},
                            {'row': 0, 'col': 2, 'text': 'std::string', 'font_name': 'Segoe UI', 'font_size': 20, 'bold': True, 'color': '#BB86FC'},
                            
                            # Row 1: Create
                            {'row': 1, 'col': 0, 'text': 'Create', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'},
                            {'row': 1, 'col': 1, 'text': 'O(1)', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'},
                            {'row': 1, 'col': 2, 'text': 'O(1)', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'},
                            
                            # Row 2: Length
                            {'row': 2, 'col': 0, 'text': 'Length', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'},
                            {'row': 2, 'col': 1, 'text': 'O(n)', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'},
                            {'row': 2, 'col': 2, 'text': 'O(1)', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'},
                            
                            # Row 3: Concatenate
                            {'row': 3, 'col': 0, 'text': 'Concatenate', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'},
                            {'row': 3, 'col': 1, 'text': 'O(n+m)', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'},
                            {'row': 3, 'col': 2, 'text': 'O(n+m)', 'font_name': 'Segoe UI', 'font_size': 16, 'color': '#E0E0E0'}
                        ]
                    },
                    {
                        'type': 'textbox',
                        'text': '"A programmer\'s evolution: First you love strings, then you hate them, then you understand them, and finally you accept that they\'ll always be a source of bugs regardless."',
                        'left': 0.5,
                        'top': 6.5,
                        'width': 12,
                        'height': 0.7,
                        'font_name': 'Segoe UI',
                        'font_size': 16,
                        'italic': True,
                        'color': '#E0E0E0',  # Light Gray
                        'alignment': 'center'
                    }
                ],
                'animations': [
                    {
                        'type': 'Fade',
                        'element': 'Title',
                        'duration': 0.7
                    },
                    {
                        'type': 'Wipe',
                        'element': 'Table header',
                        'direction': 'From Top',
                        'duration': 0.5
                    },
                    {
                        'type': 'Wipe',
                        'element': 'Table rows',
                        'direction': 'From Left',
                        'duration': 0.5,
                        'delay': 0.1  # per row
                    },
                    {
                        'type': 'Float Up',
                        'element': 'Quote',
                        'duration': 0.8,
                        'delay': 0.5
                    }
                ]
            },
            
            # String Manipulation Techniques
            {
                'layout': 'blank',
                'background': {'color': '#1A1A1A'},
                'title': {
                    'text': 'String Manipulation Techniques',
                    'left': 1,
                    'top': 0.8,
                    'width': 10,
                    'height': 1,
                    'font_name': 'Segoe UI Light',
                    'font_size': 40,
                    'bold': True,
                    'gradient': {
                        'color1': '#00FFFF',  # Cyan
                        'color2': '#00BFFF',  # Light Blue
                        'direction': 'horizontal'
                    }
                },
                'elements': [
                    {
                        'type': 'textbox',
                        'text': '1. String Traversal',
                        'left': 1,
                        'top': 2,
                        'width': 12,
                        'height': 0.7,
                        'font_name': 'Segoe UI',
                        'font_size': 28,
                        'bold': True,
                        'color': '#03DAC6'  # Light Green
                    },
                    {
                        'type': 'code',
                        'code': 'char str[] = "Hello";\nfor (int i = 0; str[i] != \'\\0\'; i++) {\n    char c = str[i];\n    // Process character c\n}',
                        'left': 1,
                        'top': 2.8,
                        'width': 5,
                        'height': 1.8,
                        'font_name': 'Consolas',
                        'font_size': 14,
                        'text_color': '#03DAC6',  # Light Green
                        'bg_color': '#252525'  # Very Dark Gray
                    },
                    {
                        'type': 'code',
                        'code': 'std::string str = "Hello";\nfor (char c : str) {\n    // Process character c\n}',
                        'left': 7,
                        'top': 2.8,
                        'width': 5,
                        'height': 1.8,
                        'font_name': 'Consolas',
                        'font_size': 14,
                        'text_color': '#03DAC6',  # Light Green
                        'bg_color': '#252525'  # Very Dark Gray
                    },
                    {
                        'type': 'textbox',
                        'text': '2. String Transformation',
                        'left': 1,
                        'top': 5,
                        'width': 12,
                        'height': 0.7,
                        'font_name': 'Segoe UI',
                        'font_size': 28,
                        'bold': True,
                        'color': '#03DAC6'  # Light Green
                    },
                    {
                        'type': 'code',
                        'code': 'std::string str = "Hello";\nstd::transform(str.begin(), str.end(),\n               str.begin(), ::toupper);\n// Result: "HELLO"',
                        'left': 1,
                        'top': 5.8,
                        'width': 5.5,
                        'height': 1.5,
                        'font_name': 'Consolas',
                        'font_size': 14,
                        'text_color': '#03DAC6',  # Light Green
                        'bg_color': '#252525'  # Very Dark Gray
                    }
                ],
                'animations': [
                    {
                        'type': 'Fade',
                        'element': 'Title',
                        'duration': 0.7
                    },
                    {
                        'type': 'Fade',
                        'element': 'String Traversal heading',
                        'duration': 0.5
                    },
                    {
                        'type': 'Appear',
                        'element': 'C-style traversal code',
                        'duration': 0.5,
                        'delay': 0.2
                    },
                    {
                        'type': 'Appear',
                        'element': 'std::string traversal code',
                        'duration': 0.5,
                        'delay': 0.3
                    },
                    {
                        'type': 'Fade',
                        'element': 'String Transformation heading',
                        'duration': 0.5,
                        'delay': 0.3
                    },
                    {
                        'type': 'Appear',
                        'element': 'transformation code',
                        'duration': 0.5,
                        'delay': 0.2
                    }
                ]
            }
        ]
    }
    
    # Write to YAML file
    with open('cpp_strings_presentation_spec.yaml', 'w') as f:
        yaml.dump(spec, f, default_flow_style=False)
    
    return spec

def main():
    # Create a sample specification file
    print("Creating demo specification file...")
    create_demo_spec()
    
    # Parse specification and generate presentation
    print("Generating presentation from specification...")
    parser = PPTSpecParser('cpp_strings_presentation_spec.yaml')
    parser.generate_presentation('CPP_Strings_Presentation.pptx')
    
    print("\nDone! The presentation has been created with notes for manual animations.")
    print("Open the presentation in PowerPoint to add the animations as described in the notes.")


if __name__ == "__main__":
    main()