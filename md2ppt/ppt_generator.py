"""
PowerPoint generator module for md2ppt.
"""

import os
from typing import List, Dict, Any, Tuple, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor

from .styler import PresentationStyler
from .utils import is_url, is_local_file


class PowerPointGenerator:
    """
    Generator for creating PowerPoint presentations from parsed markdown.
    """
    
    def __init__(self, styler: PresentationStyler):
        """
        Initialize the PowerPoint generator.
        
        Args:
            styler: PresentationStyler instance
        """
        self.styler = styler
        self.presentation = None
    
    def generate(self, slides: List[Dict[str, Any]]) -> Presentation:
        """
        Generate a PowerPoint presentation from parsed markdown.
        
        Args:
            slides: List of slide dictionaries from the parser
            
        Returns:
            A PowerPoint presentation
        """
        # Create a new presentation with the specified theme
        self.presentation = self.styler.create_presentation()
        
        # Process each slide
        for slide_data in slides:
            self._create_slide(slide_data)
        
        return self.presentation
    
    def _create_slide(self, slide_data: Dict[str, Any]) -> None:
        """
        Create a single slide in the presentation.
        
        Args:
            slide_data: Slide dictionary from the parser
        """
        # Get slide layout
        layout = self.styler.get_slide_layout(self.presentation, slide_data.get('properties', {}))
        
        # Create slide
        slide = self.presentation.slides.add_slide(layout)
        
        # Add title if present
        if slide_data.get('title'):
            if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
                self.styler.style_title(slide.shapes.title, slide_data.get('properties', {}))
        
        # Add slide content
        self._add_slide_content(slide, slide_data.get('elements', []))
        
        # Apply slide transitions and animations
        self._apply_transitions(slide, slide_data.get('properties', {}))
    
    def _add_slide_content(self, slide, elements: List[Dict[str, Any]]) -> None:
        """
        Add content elements to a slide.
        
        Args:
            slide: PowerPoint slide
            elements: List of content elements
        """
        # Starting position for content
        left = Inches(1)
        top = Inches(2) if slide.shapes.title else Inches(1)
        
        for element in elements:
            element_type = element.get('type')
            
            if element_type == 'paragraph':
                shape = self._add_text_box(slide, element['content'], left, top, 
                                          width=Inches(8), height=Inches(1))
                top += shape.height + Inches(0.2)
            
            elif element_type == 'heading':
                shape = self._add_text_box(
                    slide, element['content'], left, top,
                    width=Inches(8), height=Inches(0.8),
                    font_size=self.styler.get_heading_size(element['level']),
                    font_bold=True
                )
                top += shape.height + Inches(0.2)
            
            elif element_type in ('ordered_list', 'unordered_list'):
                shape = self._add_list(
                    slide, element['content'], left, top,
                    width=Inches(8), height=Inches(0.8),
                    is_numbered=(element_type == 'ordered_list'),
                    level=element.get('level', 1)
                )
                top += shape.height + Inches(0.2)
            
            elif element_type == 'image':
                shape = self._add_image(
                    slide, element['src'], left, top,
                    alt=element.get('alt'),
                    title=element.get('title')
                )
                if shape:
                    top += shape.height + Inches(0.2)
            
            elif element_type == 'code_block':
                shape = self._add_code_block(
                    slide, element['content'], left, top,
                    width=Inches(8), height=Inches(2),
                    language=element.get('language')
                )
                top += shape.height + Inches(0.2)
            
            elif element_type == 'quote':
                shape = self._add_quote(
                    slide, element['content'], left, top,
                    width=Inches(8), height=Inches(1.5)
                )
                top += shape.height + Inches(0.2)
    
    def _add_text_box(self, slide, text: str, left: Inches, top: Inches, 
                      width: Inches, height: Inches, font_size: Optional[Pt] = None,
                      font_bold: bool = False, font_italic: bool = False,
                      alignment: PP_ALIGN = PP_ALIGN.LEFT) -> Any:
        """
        Add a text box to a slide.
        
        Args:
            slide: PowerPoint slide
            text: Text content
            left: Left position
            top: Top position
            width: Width
            height: Height
            font_size: Font size (optional)
            font_bold: Bold text flag
            font_italic: Italic text flag
            alignment: Text alignment
            
        Returns:
            PowerPoint shape
        """
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = text
        p.alignment = alignment
        
        # Apply text formatting
        run = p.runs[0]
        if font_size:
            run.font.size = font_size
        else:
            run.font.size = self.styler.get_body_font_size()
        
        run.font.bold = font_bold
        run.font.italic = font_italic
        
        # Apply styling from the theme
        self.styler.style_body_text(run)
        
        return shape
    
    def _add_list(self, slide, items, left: Inches, top: Inches, 
                 width: Inches, height: Inches, is_numbered: bool = False,
                 level: int = 1) -> Any:
        """
        Add a list to a slide.
        
        Args:
            slide: PowerPoint slide
            items: List items
            left: Left position
            top: Top position
            width: Width
            height: Height
            is_numbered: Flag for ordered lists
            level: List level
            
        Returns:
            PowerPoint shape
        """
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        # Clear default paragraph
        first_para = text_frame.paragraphs[0]
        first_para.text = ""
        
        for i, item in enumerate(items):
            if i == 0:
                p = first_para
            else:
                p = text_frame.add_paragraph()
            
            p.text = item.get('content', '')
            p.level = level - 1  # PowerPoint levels are 0-based
            
            # Use bullet or number
            if is_numbered:
                p.font.size = self.styler.get_body_font_size()
            else:
                p.font.size = self.styler.get_body_font_size()
                
            # Apply styling
            run = p.runs[0]
            self.styler.style_body_text(run)
        
        return shape
    
    def _add_image(self, slide, src: str, left: Inches, top: Inches, 
                  alt: Optional[str] = None, title: Optional[str] = None) -> Any:
        """
        Add an image to a slide.
        
        Args:
            slide: PowerPoint slide
            src: Image source (URL or local path)
            left: Left position
            top: Top position
            alt: Alternative text
            title: Image title
            
        Returns:
            PowerPoint shape or None if image couldn't be added
        """
        try:
            # Check if it's a URL or local file
            if is_url(src):
                # Handle URL (download or process)
                import requests
                from io import BytesIO
                
                response = requests.get(src)
                image_stream = BytesIO(response.content)
                
                shape = slide.shapes.add_picture(image_stream, left, top)
            elif is_local_file(src):
                # Handle local file
                shape = slide.shapes.add_picture(src, left, top)
            else:
                # Try to handle as a local file relative to current directory
                shape = slide.shapes.add_picture(src, left, top)
            
            # Resize image if it's too large
            if shape.width > Inches(8):
                ratio = shape.height / shape.width
                shape.width = Inches(8)
                shape.height = Inches(8) * ratio
            
            # Add caption if title is provided
            if title:
                caption_top = top + shape.height + Inches(0.1)
                caption = slide.shapes.add_textbox(
                    left, caption_top, shape.width, Inches(0.5)
                )
                
                caption.text_frame.paragraphs[0].text = title
                caption.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                run = caption.text_frame.paragraphs[0].runs[0]
                run.font.italic = True
                run.font.size = Pt(10)
                
                # Adjust shape height to include caption
                shape.height += caption.height + Inches(0.1)
            
            return shape
        except Exception as e:
            print(f"Error adding image: {e}")
            # Add placeholder with error message
            shape = slide.shapes.add_textbox(left, top, Inches(4), Inches(1))
            shape.text_frame.text = f"[Image Error: {src}]"
            return shape
    
    def _add_code_block(self, slide, code: str, left: Inches, top: Inches,
                       width: Inches, height: Inches, language: Optional[str] = None) -> Any:
        """
        Add a code block to a slide.
        
        Args:
            slide: PowerPoint slide
            code: Code content
            left: Left position
            top: Top position
            width: Width
            height: Height
            language: Programming language
            
        Returns:
            PowerPoint shape
        """
        # Create a shape for the code block with background color
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        
        # Style the code block shape
        fill = rect.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background
        
        # Add border
        line = rect.line
        line.color.rgb = RGBColor(200, 200, 200)
        line.width = Pt(1)
        
        # Add text inside the shape
        text_frame = rect.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.margin_left = Pt(8)
        text_frame.margin_right = Pt(8)
        text_frame.margin_top = Pt(8)
        text_frame.margin_bottom = Pt(8)
        
        p = text_frame.paragraphs[0]
        p.text = code
        p.font.name = "Courier New"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)
        
        # If language is specified, add a label
        if language:
            lang_label = slide.shapes.add_textbox(
                left, top - Inches(0.3), Inches(2), Inches(0.2)
            )
            lang_label.text_frame.paragraphs[0].text = language
            lang_label.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
            lang_label.text_frame.paragraphs[0].runs[0].font.bold = True
            lang_label.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(100, 100, 100)
        
        return rect
    
    def _add_quote(self, slide, text: str, left: Inches, top: Inches,
                  width: Inches, height: Inches) -> Any:
        """
        Add a block quote to a slide.
        
        Args:
            slide: PowerPoint slide
            text: Quote content
            left: Left position
            top: Top position
            width: Width
            height: Height
            
        Returns:
            PowerPoint shape
        """
        # Create a textbox with indentation
        shape = slide.shapes.add_textbox(left + Inches(0.5), top, width - Inches(1), height)
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        # Style the quote
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.italic = True
        p.font.size = self.styler.get_body_font_size()
        
        # Apply quote styling
        run = p.runs[0]
        self.styler.style_quote(run)
        
        # Add a vertical line to the left (decorative)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, Inches(0.1), height
        )
        
        # Style the line
        line_fill = line.fill
        line_fill.solid()
        line_fill.fore_color.rgb = RGBColor(200, 200, 200)  # Gray line
        
        return shape
    
    def _apply_transitions(self, slide, properties: Dict[str, str]) -> None:
        """
        Apply transitions and animations to a slide.
        
        Args:
            slide: PowerPoint slide
            properties: Slide properties dictionary
        """
        # The python-pptx library has limited support for transitions and animations
        # For advanced transitions, we might need to modify the PPTX XML directly
        
        transition = properties.get('transition')
        if transition:
            # Store the transition info - will be applied later via XML modification
            if not hasattr(slide, '_md2ppt_transition'):
                slide._md2ppt_transition = {}
            
            slide._md2ppt_transition['type'] = transition
            
            # Parse additional properties like duration
            if 'duration' in properties:
                try:
                    slide._md2ppt_transition['duration'] = float(properties['duration'])
                except (ValueError, TypeError):
                    pass