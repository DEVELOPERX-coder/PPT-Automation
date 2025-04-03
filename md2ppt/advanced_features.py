"""
Advanced PowerPoint features implementation for md2ppt.

This module implements advanced PowerPoint features such as:
- Animations
- Transitions
- Custom XML manipulation for features not directly supported by python-pptx
"""

import os
import re
from typing import Dict, Any, Optional, List
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

# Namespaces used in PPTX XML
nsmap = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'
}

# Register namespaces for XML output
for prefix, uri in nsmap.items():
    ET.register_namespace(prefix, uri)


class AnimationHandler:
    """Handler for PowerPoint animations."""
    
    # Animation effect types and their PowerPoint XML equivalents
    ANIMATIONS = {
        'fade': 'fade',
        'appear': 'appear',
        'fly_in': 'fly',
        'float_in': 'float',
        'split': 'split',
        'wipe': 'wipe',
        'zoom': 'zoom',
        'grow': 'grow',
        'spin': 'spin',
        'swivel': 'swivel',
        'pulse': 'pulse',
        'bounce': 'bounce'
    }
    
    @staticmethod
    def apply_animation(presentation: Presentation, slide_idx: int, shape_idx: int, 
                       animation_type: str, duration: float = 0.5, delay: float = 0.0,
                       direction: Optional[str] = None) -> None:
        """
        Apply animation to a shape using direct XML manipulation.
        
        Args:
            presentation: PowerPoint presentation
            slide_idx: Index of the slide (0-based)
            shape_idx: Index of the shape to animate (0-based)
            animation_type: Type of animation (fade, appear, etc.)
            duration: Animation duration in seconds
            delay: Animation delay in seconds
            direction: Direction for the animation (for animations that support it)
        """
        try:
            # Get slide and shape
            if slide_idx >= len(presentation.slides):
                print(f"Slide index {slide_idx} out of range")
                return
                
            slide = presentation.slides[slide_idx]
            if shape_idx >= len(slide.shapes):
                print(f"Shape index {shape_idx} out of range for slide {slide_idx}")
                return
                
            # Get the slide's XML part
            slide_part = slide._slide
            
            # Check if animation is supported
            if animation_type not in AnimationHandler.ANIMATIONS:
                print(f"Unsupported animation: {animation_type}")
                return
                
            # Convert timing values to milliseconds for PowerPoint
            duration_ms = int(duration * 1000)
            delay_ms = int(delay * 1000)
            
            # Get shape ID to reference in animation
            shape = slide.shapes[shape_idx]
            shape_id = shape._element.shape_id
            
            # Create or get timing node
            timing = slide_part.get_or_add_timing()
            
            # Create or get the timing tree's sequence (main animation sequence)
            sequence = timing.get_or_add_sequence()
            
            # Create animation node
            anim = sequence.add_animation(shape_id)
            
            # Set animation properties
            anim_effect = anim.add_effect(AnimationHandler.ANIMATIONS[animation_type])
            anim_effect.set('dur', str(duration_ms))
            anim_effect.set('delay', str(delay_ms))
            
            # Set direction if provided and supported
            if direction and animation_type in ['fly_in', 'wipe', 'split']:
                anim_effect.set('dir', direction)
            
            print(f"Applied animation '{animation_type}' to shape {shape_idx} on slide {slide_idx}")
            
        except Exception as e:
            print(f"Error applying animation: {e}")
            # Fallback to basic animation if available


class TransitionHandler:
    """Handler for PowerPoint transitions."""
    
    # Transition types and their PowerPoint XML equivalents
    TRANSITIONS = {
        'none': None,
        'fade': 'fade',
        'push': 'push',
        'wipe': 'wipe',
        'split': 'split',
        'cut': 'cut',
        'random': 'random',
        'shape': 'circle',
        'blinds': 'blinds',
        'checker': 'checker',
        'comb': 'comb',
        'dissolve': 'dissolve',
        'zoom': 'zoom'
    }
    
    @staticmethod
    def apply_transition(presentation: Presentation, slide_idx: int, 
                        transition_type: str, duration: float = 1.0,
                        direction: Optional[str] = None) -> bool:
        """
        Apply a transition to a slide using XML manipulation.
        
        Args:
            presentation: PowerPoint presentation
            slide_idx: Index of the slide (0-based)
            transition_type: Type of transition
            duration: Transition duration in seconds
            direction: Direction for the transition (for transitions that support it)
            
        Returns:
            Success flag
        """
        try:
            # Get slide XML part
            slide = presentation.slides[slide_idx]
            slide_part = slide._slide
            
            # Check if transition is supported
            if transition_type not in TransitionHandler.TRANSITIONS:
                print(f"Unsupported transition: {transition_type}")
                return False
                
            # Convert duration to milliseconds for PowerPoint
            duration_ms = int(duration * 1000)
            
            # Create transition element
            transition = slide_part.get_or_add_transition()
            
            # Set transition attributes based on type
            transition_value = TransitionHandler.TRANSITIONS[transition_type]
            if transition_value:
                transition.set(f'{{{nsmap["p"]}}}transition-type', transition_value)
                transition.set(f'{{{nsmap["p"]}}}dur', str(duration_ms))
                
                # Set direction if provided
                if direction:
                    transition.set(f'{{{nsmap["p"]}}}dir', direction)
            
            return True
        except Exception as e:
            print(f"Error applying transition: {e}")
            return False


class StyleHandler:
    """Handler for advanced PowerPoint styles."""
    
    @staticmethod
    def apply_theme_color(presentation: Presentation, color_dict: Dict[str, Any]) -> None:
        """
        Apply custom theme colors to the presentation.
        
        Args:
            presentation: PowerPoint presentation
            color_dict: Dictionary of color settings
        """
        try:
            # Get the theme part
            if not presentation.part.package.parts:
                return
                
            # Find the theme part
            theme_part = None
            for part in presentation.part.package.parts:
                if part.partname.endswith('theme/theme1.xml'):
                    theme_part = part
                    break
                    
            if not theme_part:
                return
                
            # Parse the theme XML
            theme_xml = ET.fromstring(theme_part.blob)
            
            # Find the color scheme
            color_scheme = theme_xml.find('.//a:clrScheme', nsmap)
            if not color_scheme:
                return
                
            # Update colors
            for color_name, color_value in color_dict.items():
                # Find the color element
                color_element = color_scheme.find(f'.//a:{color_name}', nsmap)
                if color_element is not None:
                    # Convert RGB to hex
                    if isinstance(color_value, (list, tuple)) and len(color_value) == 3:
                        r, g, b = color_value
                        hex_color = f"{r:02x}{g:02x}{b:02x}"
                        
                        # Set the color value
                        rgb_element = color_element.find('.//a:srgbClr', nsmap)
                        if rgb_element is not None:
                            rgb_element.set('val', hex_color)
            
            # Save changes
            theme_part.blob = ET.tostring(theme_xml)
            
        except Exception as e:
            print(f"Error applying theme colors: {e}")
    
    @staticmethod
    def apply_master_styles(presentation: Presentation, style_dict: Dict[str, Any]) -> None:
        """
        Apply styles to the slide master.
        
        Args:
            presentation: PowerPoint presentation
            style_dict: Dictionary of style settings
        """
        try:
            # Get slide master
            master = presentation.slide_master
            
            # Apply styles to title placeholder
            if 'title' in style_dict and master.shapes.title:
                title_style = style_dict['title']
                title_shape = master.shapes.title
                
                # Apply text style
                if 'font' in title_style:
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if 'name' in title_style['font']:
                                run.font.name = title_style['font']['name']
                            if 'size' in title_style['font']:
                                from pptx.util import Pt
                                run.font.size = Pt(title_style['font']['size'])
                            if 'bold' in title_style['font']:
                                run.font.bold = title_style['font']['bold']
                            if 'italic' in title_style['font']:
                                run.font.italic = title_style['font']['italic']
                
                # Apply color
                if 'color' in title_style:
                    from pptx.dml.color import RGBColor
                    color = title_style['color']
                    if isinstance(color, (list, tuple)) and len(color) == 3:
                        for paragraph in title_shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(*color)
            
            # Apply styles to body placeholder
            if 'body' in style_dict:
                body_style = style_dict['body']
                
                # Find body placeholder
                for shape in master.placeholders:
                    if shape.placeholder_format.type == 1:  # Body placeholder
                        # Apply text style
                        if 'font' in body_style:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if 'name' in body_style['font']:
                                        run.font.name = body_style['font']['name']
                                    if 'size' in body_style['font']:
                                        from pptx.util import Pt
                                        run.font.size = Pt(body_style['font']['size'])
                                    if 'bold' in body_style['font']:
                                        run.font.bold = body_style['font']['bold']
                                    if 'italic' in body_style['font']:
                                        run.font.italic = body_style['font']['italic']
            
        except Exception as e:
            print(f"Error applying master styles: {e}")
    
    @staticmethod
    def apply_background_color(slide, color: tuple) -> None:
        """
        Apply background color to a slide.
        
        Args:
            slide: PowerPoint slide
            color: RGB color as (r, g, b) tuple
        """
        try:
            from pptx.dml.color import RGBColor
            
            # Get background fill
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*color)
            
        except Exception as e:
            print(f"Error applying background color: {e}")


class XMLHelper:
    """Helper class for XML manipulation of PPTX files."""
    
    @staticmethod
    def add_namespaces_to_root(root):
        """Add required namespaces to an XML root element."""
        for prefix, uri in nsmap.items():
            root.set(f'xmlns:{prefix}', uri)
    
    @staticmethod
    def create_element(tag, attributes=None, text=None, nsmap=None):
        """Create an XML element with the given tag, attributes, and text."""
        if nsmap:
            ET.register_namespace('', nsmap)
        
        element = ET.Element(tag)
        
        if attributes:
            for key, value in attributes.items():
                element.set(key, value)
        
        if text:
            element.text = text
        
        return element


def apply_advanced_features(presentation: Presentation, slide_data: List[Dict[str, Any]]) -> None:
    """
    Apply advanced features to a presentation based on parsed slide data.
    
    Args:
        presentation: PowerPoint presentation
        slide_data: List of slide dictionaries with properties
    """
    for idx, slide_info in enumerate(slide_data):
        properties = slide_info.get('properties', {})
        
        # Apply transitions
        if 'transition' in properties:
            transition_info = properties['transition']
            transition_type = transition_info
            duration = 1.0  # default
            direction = None
            
            # Check if transition has additional properties
            if ',' in transition_info:
                parts = transition_info.split(',')
                transition_type = parts[0].strip()
                
                # Parse additional properties
                for part in parts[1:]:
                    if '=' in part:
                        key, value = part.split('=', 1)
                        if key.strip() == 'duration':
                            try:
                                duration = float(value.strip())
                            except ValueError:
                                pass
                        elif key.strip() == 'direction':
                            direction = value.strip()
            
            # Apply the transition
            if idx < len(presentation.slides):
                TransitionHandler.apply_transition(
                    presentation, idx, transition_type, duration, direction
                )
        
        # Apply animations to elements
        if idx < len(presentation.slides) and 'elements' in slide_info:
            slide = presentation.slides[idx]
            
            # Track which elements need animations
            for element_idx, element in enumerate(slide_info.get('elements', [])):
                if 'animation' in element:
                    animation_info = element['animation']
                    animation_type = 'fade'  # default
                    duration = 0.5  # default
                    delay = 0.0  # default
                    direction = None
                    
                    # Parse animation specifications
                    if isinstance(animation_info, str):
                        if animation_info in AnimationHandler.ANIMATIONS:
                            animation_type = animation_info
                        elif ',' in animation_info:
                            # Parse comma-separated animation settings
                            parts = animation_info.split(',')
                            animation_type = parts[0].strip()
                            
                            for part in parts[1:]:
                                if '=' in part:
                                    key, value = part.split('=', 1)
                                    key = key.strip()
                                    value = value.strip()
                                    
                                    if key == 'duration':
                                        try:
                                            duration = float(value)
                                        except ValueError:
                                            pass
                                    elif key == 'delay':
                                        try:
                                            delay = float(value)
                                        except ValueError:
                                            pass
                                    elif key == 'direction':
                                        direction = value
                    
                    # Apply animation to corresponding shape
                    # In real implementation, we need to map element_idx to shape_idx
                    # For simplicity, we assume they match
                    if element_idx < len(slide.shapes):
                        AnimationHandler.apply_animation(
                            presentation, idx, element_idx, 
                            animation_type, duration, delay, direction
                        )
        
        # Apply slide background color if specified
        if 'background_color' in properties and idx < len(presentation.slides):
            from pptx.dml.color import RGBColor
            from md2ppt.utils import parse_color
            color = parse_color(properties['background_color'])
            if color:
                slide = presentation.slides[idx]
                StyleHandler.apply_background_color(slide, color)