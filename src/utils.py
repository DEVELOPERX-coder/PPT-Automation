"""
Utility functions for the PowerPoint Generator.

This module provides helper functions for the PowerPoint Generator.
"""

import os
import re
import logging
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

def format_text_frame(text_frame, font=None, size=None, color=None, bold=None, 
                     italic=None, underline=None, alignment=None, line_spacing=None):
    """
    Apply formatting to a text frame.
    
    Args:
        text_frame: The text frame to format.
        font (str, optional): Font name.
        size (int, optional): Font size in points.
        color (tuple or list, optional): RGB color as (r, g, b) tuple.
        bold (bool, optional): Whether to make text bold.
        italic (bool, optional): Whether to make text italic.
        underline (bool, optional): Whether to underline text.
        alignment (PP_ALIGN, optional): Text alignment.
        line_spacing (float, optional): Line spacing multiplier.
    """
    if not hasattr(text_frame, 'paragraphs'):
        return
    
    # Set paragraph properties that apply to all runs
    if alignment is not None:
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = alignment
    
    if line_spacing is not None:
        for paragraph in text_frame.paragraphs:
            paragraph.line_spacing = line_spacing
    
    # Set run-level properties
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if font is not None:
                run.font.name = font
            
            if size is not None:
                run.font.size = Pt(size)
            
            if color is not None:
                if isinstance(color, (tuple, list)) and len(color) == 3:
                    run.font.color.rgb = RGBColor(*color)
            
            if bold is not None:
                run.font.bold = bold
            
            if italic is not None:
                run.font.italic = italic
            
            if underline is not None:
                run.font.underline = underline

def apply_theme_settings(presentation, theme_settings):
    """
    Apply theme settings to a presentation.
    
    Args:
        presentation: The presentation to apply theme settings to.
        theme_settings (dict): Dictionary of theme settings.
    """
    # Note: This function is a placeholder for applying theme settings globally.
    # python-pptx has limited support for applying global theme settings.
    # Most formatting is applied at the individual element level.
    logger.debug("Applied theme settings to presentation")

def resolve_variables(data, variables):
    """
    Resolve variables in a data structure.
    
    Args:
        data: The data structure to process (dict, list, str).
        variables (dict): Dictionary of variables.
        
    Returns:
        The processed data structure with variables resolved.
    """
    if isinstance(data, dict):
        return {k: resolve_variables(v, variables) for k, v in data.items()}
    elif isinstance(data, list):
        return [resolve_variables(item, variables) for item in data]
    elif isinstance(data, str):
        # Replace variables in string
        result = data
        for var_name, var_value in variables.items():
            # Match both {{var}} and ${var} syntax
            result = result.replace('{{' + var_name + '}}', str(var_value))
            result = result.replace('${' + var_name + '}', str(var_value))
        return result
    else:
        return data

def get_rgb_color(color_value):
    """
    Parse a color value into an RGB tuple.
    
    Args:
        color_value: Color value to parse (string, list, or tuple).
        
    Returns:
        tuple: RGB color as (r, g, b) tuple.
    """
    print(f"Input type: {type(color_value)}, value: {color_value}")

    if isinstance(color_value, (list, tuple)) and len(color_value) == 3:
        return tuple(int(c) for c in color_value)
    
    elif isinstance(color_value, str):
        # Handle hex color
        if color_value.startswith('#'):
            hex_color = color_value.lstrip('#')
            if len(hex_color) == 3:
                hex_color = ''.join([c+c for c in hex_color])
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        
        # Handle named colors
        named_colors = {
            'black': (0, 0, 0),
            'white': (255, 255, 255),
            'red': (255, 0, 0),
            'green': (0, 128, 0),
            'blue': (0, 0, 255),
            'yellow': (255, 255, 0),
            'purple': (128, 0, 128),
            'orange': (255, 165, 0),
            'gray': (128, 128, 128),
            'light_gray': (211, 211, 211),
            'dark_gray': (64, 64, 64),
            'cyan': (0, 255, 255),
            'magenta': (255, 0, 255),
            'pink': (255, 192, 203),
            'brown': (165, 42, 42),
            'navy': (0, 0, 128),
            'teal': (0, 128, 128)
        }
        
        color_lower = color_value.lower().replace(' ', '_')
        if color_lower in named_colors:
            return named_colors[color_lower]
    
    # Default to black if parsing fails
    logger.warning(f"Could not parse color value: {color_value}, using black")
    return (0, 0, 0)

def get_image_dimensions(image_path):
    """
    Get the dimensions of an image file.
    
    Args:
        image_path (str): Path to the image file.
        
    Returns:
        tuple: (width, height) in pixels, or None if the file doesn't exist.
    """
    try:
        from PIL import Image
        
        if not os.path.exists(image_path):
            logger.warning(f"Image file not found: {image_path}")
            return None
        
        with Image.open(image_path) as img:
            return img.size
    except Exception as e:
        logger.error(f"Error getting image dimensions: {e}")
        return None

def calculate_aspect_ratio(width, height, max_width=None, max_height=None):
    """
    Calculate dimensions while preserving aspect ratio.
    
    Args:
        width (float): Original width.
        height (float): Original height.
        max_width (float, optional): Maximum allowed width.
        max_height (float, optional): Maximum allowed height.
        
    Returns:
        tuple: (new_width, new_height) that preserves the aspect ratio.
    """
    aspect_ratio = width / height
    
    if max_width and max_height:
        if width > max_width or height > max_height:
            width_ratio = max_width / width
            height_ratio = max_height / height
            
            # Use the smaller ratio to ensure both dimensions fit within limits
            if width_ratio < height_ratio:
                return (max_width, max_width / aspect_ratio)
            else:
                return (max_height * aspect_ratio, max_height)
    
    elif max_width and width > max_width:
        return (max_width, max_width / aspect_ratio)
    
    elif max_height and height > max_height:
        return (max_height * aspect_ratio, max_height)
    
    return (width, height)

def sanitize_filename(filename):
    """
    Sanitize a filename by removing invalid characters.
    
    Args:
        filename (str): The filename to sanitize.
        
    Returns:
        str: Sanitized filename.
    """
    # Remove invalid characters
    invalid_chars = r'[<>:"/\\|?*]'
    sanitized = re.sub(invalid_chars, '_', filename)
    
    # Ensure the filename isn't too long (Windows has a 255 character limit)
    if len(sanitized) > 240:
        base, ext = os.path.splitext(sanitized)
        sanitized = base[:240 - len(ext)] + ext
    
    return sanitized