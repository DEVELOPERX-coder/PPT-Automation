"""
Slide Builder Module

This module is responsible for creating and configuring individual slides
based on YAML configuration.
"""

import logging
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from src.element_factory import ElementFactory
from src.utils import format_text_frame

logger = logging.getLogger(__name__)

class SlideBuilder:
    """
    A class for building individual slides in a PowerPoint presentation.
    """
    
    def __init__(self, presentation):
        """
        Initialize the SlideBuilder with a presentation.
        
        Args:
            presentation: The pptx.Presentation object to add slides to.
        """
        self.presentation = presentation
        self.element_factory = ElementFactory()
    
    def create_slide(self, slide_data, theme_settings):
        """
        Create a slide based on the provided data.
        
        Args:
            slide_data (dict): Dictionary containing slide data.
            theme_settings (dict): Dictionary of theme settings.
            
        Returns:
            Slide: The created slide object.
        """
        # Get slide type
        slide_type = slide_data.get('type', 'blank')
        
        # Get the appropriate layout
        layout = self._get_slide_layout(slide_type)
        
        # Create the slide
        slide = self.presentation.slides.add_slide(layout)
        
        # Apply background
        self._apply_background(slide, slide_data, theme_settings)
        
        # Process slide content based on type
        if slide_type == 'title':
            self._create_title_slide(slide, slide_data, theme_settings)
        elif slide_type == 'title_and_content':
            self._create_title_content_slide(slide, slide_data, theme_settings)
        elif slide_type == 'section':
            self._create_section_slide(slide, slide_data, theme_settings)
        elif slide_type == 'two_content':
            self._create_two_content_slide(slide, slide_data, theme_settings)
        elif slide_type == 'title_only':
            self._create_title_only_slide(slide, slide_data, theme_settings)
        elif slide_type == 'blank':
            self._create_blank_slide(slide, slide_data, theme_settings)
        else:
            logger.warning(f"Unknown slide type: {slide_type}, defaulting to blank")
            self._create_blank_slide(slide, slide_data, theme_settings)
        
        # Add custom elements
        if 'elements' in slide_data:
            self._add_custom_elements(slide, slide_data['elements'], theme_settings)
        
        # Apply animations if specified
        if 'animations' in slide_data:
            self._apply_animations(slide, slide_data['animations'])
        
        return slide
    
    def _get_slide_layout(self, slide_type):
        """
        Get the appropriate slide layout for the given slide type.
        
        Args:
            slide_type (str): Type of slide to create.
            
        Returns:
            SlideLayout: The slide layout to use.
        """
        # Map slide types to indices in the slide layout collection
        slide_layouts = {
            'title': 0,              # Title Slide
            'title_and_content': 1,  # Title and Content
            'section': 2,            # Section Header
            'two_content': 3,        # Two Content
            'comparison': 4,         # Comparison
            'title_only': 5,         # Title Only
            'blank': 6,              # Blank
            'content_with_caption': 7,  # Content with Caption
            'picture_with_caption': 8   # Picture with Caption
        }
        
        layout_idx = slide_layouts.get(slide_type, 6)  # Default to blank
        
        # Handle case where presentation doesn't have enough layouts
        if layout_idx >= len(self.presentation.slide_layouts):
            logger.warning(f"Slide layout index {layout_idx} not available, using blank (6)")
            layout_idx = min(6, len(self.presentation.slide_layouts) - 1)
        
        return self.presentation.slide_layouts[layout_idx]
    
    def _apply_background(self, slide, slide_data, theme_settings):
        """
        Apply background to a slide.
        
        Args:
            slide: The slide to modify.
            slide_data (dict): Slide configuration data.
            theme_settings (dict): Theme settings.
        """
        # Check if slide has specific background
        if 'background' in slide_data:
            bg_data = slide_data['background']
            
            # Handle solid color background
            if 'color' in bg_data:
                color_value = bg_data['color']
                # Convert color value to RGB tuple
                if isinstance(color_value, str) and color_value.startswith('#'):
                    hex_color = color_value.lstrip('#')
                    if len(hex_color) == 3:
                        hex_color = ''.join([c+c for c in hex_color])
                    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                elif isinstance(color_value, (list, tuple)) and len(color_value) == 3:
                    rgb = tuple(color_value)
                else:
                    rgb = theme_settings['background_color']
                
                # Apply the background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*rgb)
            
            # Handle image background
            elif 'image' in bg_data:
                image_path = bg_data['image']
                if os.path.exists(image_path):
                    # Note: As of my knowledge cutoff, python-pptx doesn't support
                    # setting image as slide background directly.
                    # As a workaround, we'll add an image that covers the entire slide
                    left = Inches(0)
                    top = Inches(0)
                    width = self.presentation.slide_width
                    height = self.presentation.slide_height
                    slide.shapes.add_picture(image_path, left, top, width, height)
                    logger.debug(f"Added image background: {image_path}")
                else:
                    logger.warning(f"Background image not found: {image_path}")
        else:
            # Use theme background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*theme_settings['background_color'])
    
    def _create_title_slide(self, slide, slide_data, theme_settings):
        """
        Create a title slide.
        
        Args:
            slide: The slide to populate.
            slide_data (dict): Data for the slide.
            theme_settings (dict): Theme settings.
        """
        if 'title' in slide_data and hasattr(slide.shapes, 'title'):
            title = slide.shapes.title
            title.text = slide_data['title']
            format_text_frame(
                title.text_frame,
                font=theme_settings['title_font'],
                size=theme_settings['title_font_size'],
                color=theme_settings['title_color']
            )
        
        if 'subtitle' in slide_data:
            # Find subtitle placeholder
            subtitle = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # Subtitle
                    subtitle = shape
                    break
            
            if subtitle:
                subtitle.text = slide_data['subtitle']
                format_text_frame(
                    subtitle.text_frame,
                    font=theme_settings['subtitle_font'],
                    size=theme_settings['subtitle_font_size'],
                    color=theme_settings['text_color']
                )
    
    def _create_title_content_slide(self, slide, slide_data, theme_settings):
        """
        Create a title and content slide.
        
        Args:
            slide: The slide to populate.
            slide_data (dict): Data for the slide.
            theme_settings (dict): Theme settings.
        """
        # Add title
        if 'title' in slide_data and hasattr(slide.shapes, 'title'):
            title = slide.shapes.title
            title.text = slide_data['title']
            format_text_frame(
                title.text_frame,
                font=theme_settings['title_font'],
                size=theme_settings['title_font_size'],
                color=theme_settings['title_color']
            )
        
        # Add content
        if 'content' in slide_data:
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Content
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                content_data = slide_data['content']
                
                # Handle different content types
                if isinstance(content_data, str):
                    # Plain text content
                    content_placeholder.text = content_data
                    format_text_frame(
                        content_placeholder.text_frame,
                        font=theme_settings['body_font'],
                        size=theme_settings['body_font_size'],
                        color=theme_settings['text_color']
                    )
                
                elif isinstance(content_data, list):
                    # Bullet point list
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for item in content_data:
                        p = text_frame.add_paragraph()
                        p.text = item
                        p.level = 0  # Top level bullet
                        
                        # Format the paragraph
                        p.font.name = theme_settings['body_font']
                        p.font.size = Pt(theme_settings['body_font_size'])
                        p.font.color.rgb = RGBColor(*theme_settings['text_color'])
                
                elif isinstance(content_data, dict) and 'type' in content_data:
                    # It's a complex content element (table, chart, etc.)
                    content_type = content_data['type']
                    
                    if content_type == 'table':
                        self.element_factory.create_table(
                            slide, 
                            content_placeholder.left,
                            content_placeholder.top,
                            content_placeholder.width,
                            content_placeholder.height,
                            content_data, 
                            theme_settings
                        )
                    
                    elif content_type == 'chart':
                        self.element_factory.create_chart(
                            slide, 
                            content_placeholder.left,
                            content_placeholder.top,
                            content_placeholder.width,
                            content_placeholder.height,
                            content_data, 
                            theme_settings
                        )
                    
                    elif content_type == 'image':
                        self.element_factory.create_image(
                            slide, 
                            content_placeholder.left,
                            content_placeholder.top,
                            content_placeholder.width,
                            content_placeholder.height,
                            content_data
                        )
                    
                    elif content_type == 'code':
                        self.element_factory.create_code_block(
                            slide, 
                            content_placeholder.left,
                            content_placeholder.top,
                            content_placeholder.width,
                            content_placeholder.height,
                            content_data, 
                            theme_settings
                        )
    
    def _create_section_slide(self, slide, slide_data, theme_settings):
        """
        Create a section header slide.
        
        Args:
            slide: The slide to populate.
            slide_data (dict): Data for the slide.
            theme_settings (dict): Theme settings.
        """
        if 'title' in slide_data and hasattr(slide.shapes, 'title'):
            title = slide.shapes.title
            title.text = slide_data['title']
            format_text_frame(
                title.text_frame,
                font=theme_settings['title_font'],
                size=theme_settings['title_font_size'],
                color=theme_settings['accent_color'],
                bold=True,
                alignment=PP_ALIGN.CENTER
            )
    
    def _create_two_content_slide(self, slide, slide_data, theme_settings):
        """
        Create a slide with two content columns.
        
        Args:
            slide: The slide to populate.
            slide_data (dict): Data for the slide.
            theme_settings (dict): Theme settings.
        """
        # Add title
        if 'title' in slide_data and hasattr(slide.shapes, 'title'):
            title = slide.shapes.title
            title.text = slide_data['title']
            format_text_frame(
                title.text_frame,
                font=theme_settings['title_font'],
                size=theme_settings['title_font_size'],
                color=theme_settings['title_color']
            )
        
        # Find left and right placeholders
        left_placeholder = None
        right_placeholder = None
        
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 7:  # Content
                if not left_placeholder:
                    left_placeholder = shape
                else:
                    right_placeholder = shape
                    break
        
        # Add left content
        if 'left_content' in slide_data and left_placeholder:
            self._add_content_to_placeholder(
                left_placeholder,
                slide_data['left_content'],
                slide,
                theme_settings
            )
        
        # Add right content
        if 'right_content' in slide_data and right_placeholder:
            self._add_content_to_placeholder(
                right_placeholder,
                slide_data['right_content'],
                slide,
                theme_settings
            )
    
    def _create_title_only_slide(self, slide, slide_data, theme_settings):
        """
        Create a title only slide.
        
        Args:
            slide: The slide to populate.
            slide_data (dict): Data for the slide.
            theme_settings (dict): Theme settings.
        """
        if 'title' in slide_data and hasattr(slide.shapes, 'title'):
            title = slide.shapes.title
            title.text = slide_data['title']
            format_text_frame(
                title.text_frame,
                font=theme_settings['title_font'],
                size=theme_settings['title_font_size'],
                color=theme_settings['title_color']
            )
    
    def _create_blank_slide(self, slide, slide_data, theme_settings):
        """
        Create a blank slide.
        
        Args:
            slide: The slide to populate.
            slide_data (dict): Data for the slide.
            theme_settings (dict): Theme settings.
        """
        # Blank slides have no default content, only custom elements
        pass
    
    def _add_content_to_placeholder(self, placeholder, content_data, slide, theme_settings):
        """
        Add content to a placeholder based on its type.
        
        Args:
            placeholder: The placeholder shape.
            content_data: The content to add.
            slide: The slide containing the placeholder.
            theme_settings (dict): Theme settings.
        """
        if isinstance(content_data, str):
            # Plain text content
            placeholder.text = content_data
            format_text_frame(
                placeholder.text_frame,
                font=theme_settings['body_font'],
                size=theme_settings['body_font_size'],
                color=theme_settings['text_color']
            )
        
        elif isinstance(content_data, list):
            # Bullet point list
            text_frame = placeholder.text_frame
            text_frame.clear()
            
            for item in content_data:
                p = text_frame.add_paragraph()
                p.text = str(item)
                p.level = 0  # Top level bullet
                
                # Format the paragraph
                p.font.name = theme_settings['body_font']
                p.font.size = Pt(theme_settings['body_font_size'])
                p.font.color.rgb = RGBColor(*theme_settings['text_color'])
        
        elif isinstance(content_data, dict) and 'type' in content_data:
            # It's a complex content element (table, chart, etc.)
            content_type = content_data['type']
            
            if content_type == 'table':
                self.element_factory.create_table(
                    slide, 
                    placeholder.left,
                    placeholder.top,
                    placeholder.width,
                    placeholder.height,
                    content_data, 
                    theme_settings
                )
            
            elif content_type == 'chart':
                self.element_factory.create_chart(
                    slide, 
                    placeholder.left,
                    placeholder.top,
                    placeholder.width,
                    placeholder.height,
                    content_data, 
                    theme_settings
                )
            
            elif content_type == 'image':
                self.element_factory.create_image(
                    slide, 
                    placeholder.left,
                    placeholder.top,
                    placeholder.width,
                    placeholder.height,
                    content_data
                )
            
            elif content_type == 'code':
                self.element_factory.create_code_block(
                    slide, 
                    placeholder.left,
                    placeholder.top,
                    placeholder.width,
                    placeholder.height,
                    content_data, 
                    theme_settings
                )
    
    def _add_custom_elements(self, slide, elements_data, theme_settings):
        """
        Add custom elements to a slide.
        
        Args:
            slide: The slide to add elements to.
            elements_data (list): List of element data dictionaries.
            theme_settings (dict): Theme settings.
        """
        for element_data in elements_data:
            element_type = element_data.get('type', '')
            
            if element_type == 'text_box':
                self.element_factory.create_text_box(slide, element_data, theme_settings)
            
            elif element_type == 'shape':
                self.element_factory.create_shape(slide, element_data, theme_settings)
            
            elif element_type == 'image':
                self.element_factory.create_image(slide, None, None, None, None, element_data)
            
            elif element_type == 'table':
                self.element_factory.create_table(slide, None, None, None, None, element_data, theme_settings)
            
            elif element_type == 'chart':
                self.element_factory.create_chart(slide, None, None, None, None, element_data, theme_settings)
            
            elif element_type == 'code':
                self.element_factory.create_code_block(slide, None, None, None, None, element_data, theme_settings)
            
            else:
                logger.warning(f"Unknown element type: {element_type}")
    
    def _apply_animations(self, slide, animations_data):
        """
        Apply animations to slide elements.
        
        Args:
            slide: The slide to apply animations to.
            animations_data (list): List of animation configurations.
        """
        # Note: As of my knowledge cutoff, python-pptx doesn't support
        # setting animations programmatically. This is a placeholder for
        # future implementation if the library adds support.
        logger.warning("Animations are not currently supported by python-pptx")