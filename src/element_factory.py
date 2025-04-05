"""
Element Factory Module

This module is responsible for creating individual PowerPoint elements
such as shapes, text boxes, tables, and charts.
"""

import os
import logging
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from src.utils import format_text_frame, get_rgb_color

logger = logging.getLogger(__name__)

class ElementFactory:
    """
    A factory class for creating PowerPoint elements.
    """
    
    def create_text_box(self, slide, element_data, theme_settings):
        """
        Create a text box element.
        
        Args:
            slide: The slide to add the text box to.
            element_data (dict): Text box configuration data.
            theme_settings (dict): Theme settings.
            
        Returns:
            The created text box shape.
        """
        # Get position and size
        left = Inches(element_data.get('left', 1))
        top = Inches(element_data.get('top', 1))
        width = Inches(element_data.get('width', 4))
        height = Inches(element_data.get('height', 1))
        
        # Create text box
        text_box = slide.shapes.add_textbox(left, top, width, height)
        
        # Set text content
        if 'text' in element_data:
            text_box.text = element_data['text']
        
        # Format text
        font = element_data.get('font', theme_settings['body_font'])
        size = element_data.get('size', theme_settings['body_font_size'])
        
        # Get color
        if 'color' in element_data:
            color = get_rgb_color(element_data['color'])
        else:
            color = theme_settings['text_color']
        
        # Get alignment
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        alignment = align_map.get(element_data.get('align', 'left'), PP_ALIGN.LEFT)
        
        # Apply formatting
        format_text_frame(
            text_box.text_frame,
            font=font,
            size=size,
            color=color,
            alignment=alignment,
            bold=element_data.get('bold', False),
            italic=element_data.get('italic', False),
            underline=element_data.get('underline', False)
        )
        
        return text_box
    
    def create_shape(self, slide, element_data, theme_settings):
        """
        Create a shape element.
        
        Args:
            slide: The slide to add the shape to.
            element_data (dict): Shape configuration data.
            theme_settings (dict): Theme settings.
            
        Returns:
            The created shape.
        """
        # Get position and size
        left = Inches(element_data.get('left', 1))
        top = Inches(element_data.get('top', 1))
        width = Inches(element_data.get('width', 2))
        height = Inches(element_data.get('height', 1))
        
        # Get shape type
        shape_type = element_data.get('shape_type', 'rectangle').lower()
        
        # Map shape names to MSO_AUTO_SHAPE_TYPE
        shape_map = {
            'rectangle': MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            'rounded_rectangle': MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            'oval': MSO_AUTO_SHAPE_TYPE.OVAL,
            'diamond': MSO_AUTO_SHAPE_TYPE.DIAMOND,
            'triangle': MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
            'right_triangle': MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
            'pentagon': MSO_AUTO_SHAPE_TYPE.PENTAGON,
            'hexagon': MSO_AUTO_SHAPE_TYPE.HEXAGON,
            'heptagon': MSO_AUTO_SHAPE_TYPE.HEPTAGON,
            'octagon': MSO_AUTO_SHAPE_TYPE.OCTAGON,
            'star': MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
            'star4': MSO_AUTO_SHAPE_TYPE.STAR_4_POINT,
            'star5': MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
            'star6': MSO_AUTO_SHAPE_TYPE.STAR_6_POINT,
            'star7': MSO_AUTO_SHAPE_TYPE.STAR_7_POINT,
            'star8': MSO_AUTO_SHAPE_TYPE.STAR_8_POINT,
            'arrow': MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            'up_arrow': MSO_AUTO_SHAPE_TYPE.UP_ARROW,
            'down_arrow': MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
            'left_arrow': MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
            'right_arrow': MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            'left_right_arrow': MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW,
            'up_down_arrow': MSO_AUTO_SHAPE_TYPE.UP_DOWN_ARROW,
            'cloud': MSO_AUTO_SHAPE_TYPE.CLOUD,
            'heart': MSO_AUTO_SHAPE_TYPE.HEART,
            'lightning': MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT,
            'sun': MSO_AUTO_SHAPE_TYPE.SUN,
            'moon': MSO_AUTO_SHAPE_TYPE.MOON,
            'smiley': MSO_AUTO_SHAPE_TYPE.SMILEY_FACE,
            'no_symbol': MSO_AUTO_SHAPE_TYPE.NO_SYMBOL,
            'arc': MSO_AUTO_SHAPE_TYPE.ARC,
            'plaque': MSO_AUTO_SHAPE_TYPE.PLAQUE,
            'can': MSO_AUTO_SHAPE_TYPE.CAN,
            'cube': MSO_AUTO_SHAPE_TYPE.CUBE,
            'bevel': MSO_AUTO_SHAPE_TYPE.BEVEL,
            'donut': MSO_AUTO_SHAPE_TYPE.DONUT,
            'pie': MSO_AUTO_SHAPE_TYPE.PIE,
            'block_arc': MSO_AUTO_SHAPE_TYPE.BLOCK_ARC,
            'folded_corner': MSO_AUTO_SHAPE_TYPE.FOLDED_CORNER,
            'frame': MSO_AUTO_SHAPE_TYPE.FRAME
        }
        
        shape_enum = shape_map.get(shape_type, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        
        # Create shape
        shape = slide.shapes.add_shape(shape_enum, left, top, width, height)
        
        # Set shape fill
        if 'fill_color' in element_data:
            fill_color = get_rgb_color(element_data['fill_color'])
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)
        
        # Set shape line
        if 'line_color' in element_data:
            line_color = get_rgb_color(element_data['line_color'])
            shape.line.color.rgb = RGBColor(*line_color)
        
        # Set line width if specified
        if 'line_width' in element_data:
            shape.line.width = Pt(element_data['line_width'])
        
        # Add text if specified
        if 'text' in element_data:
            shape.text = element_data['text']
            
            # Format text
            font = element_data.get('font', theme_settings['body_font'])
            size = element_data.get('size', theme_settings['body_font_size'])
            
            # Get color
            if 'text_color' in element_data:
                color = get_rgb_color(element_data['text_color'])
            else:
                color = theme_settings['text_color']
            
            # Get alignment
            align_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT,
                'justify': PP_ALIGN.JUSTIFY
            }
            alignment = align_map.get(element_data.get('align', 'center'), PP_ALIGN.CENTER)
            
            # Apply formatting
            format_text_frame(
                shape.text_frame,
                font=font,
                size=size,
                color=color,
                alignment=alignment,
                bold=element_data.get('bold', False),
                italic=element_data.get('italic', False),
                underline=element_data.get('underline', False)
            )
        
        return shape
    
    def create_image(self, slide, left, top, width, height, element_data):
        """
        Create an image element.
        
        Args:
            slide: The slide to add the image to.
            left: Left position (can be None if specified in element_data).
            top: Top position (can be None if specified in element_data).
            width: Width (can be None if specified in element_data).
            height: Height (can be None if specified in element_data).
            element_data (dict): Image configuration data.
            
        Returns:
            The created image shape or None if image not found.
        """
        # Get image path
        if 'path' not in element_data:
            logger.error("Image path not specified")
            return None
        
        image_path = element_data['path']
        if not os.path.exists(image_path):
            logger.error(f"Image file not found: {image_path}")
            return None
        
        # Get position and size
        left = left or Inches(element_data.get('left', 1))
        top = top or Inches(element_data.get('top', 1))
        
        # If width and height are not specified, use the image's natural size
        if ('width' not in element_data and width is None) or ('height' not in element_data and height is None):
            image_shape = slide.shapes.add_picture(image_path, left, top)
            
            # Scale the image if only one dimension is specified
            if 'width' in element_data and width is None:
                width = Inches(element_data['width'])
                scale_factor = width / image_shape.width
                height = image_shape.height * scale_factor
            elif 'height' in element_data and height is None:
                height = Inches(element_data['height'])
                scale_factor = height / image_shape.height
                width = image_shape.width * scale_factor
            
            # Remove the temporary image shape
            sp_id = image_shape.element.attrib.get('id')
            if sp_id:
                old_elm = image_shape._element
                old_elm.getparent().remove(old_elm)
        else:
            width = width or Inches(element_data.get('width', 4))
            height = height or Inches(element_data.get('height', 3))
        
        # Create the image with the determined dimensions
        image_shape = slide.shapes.add_picture(image_path, left, top, width, height)
        
        return image_shape
    
    def create_table(self, slide, left, top, width, height, element_data, theme_settings):
        """
        Create a table element.
        
        Args:
            slide: The slide to add the table to.
            left: Left position (can be None if specified in element_data).
            top: Top position (can be None if specified in element_data).
            width: Width (can be None).
            height: Height (can be None).
            element_data (dict): Table configuration data.
            theme_settings (dict): Theme settings.
            
        Returns:
            The created table shape.
        """
        # Get table data
        if 'data' not in element_data or not element_data['data']:
            logger.error("Table data not specified or empty")
            return None
        
        table_data = element_data['data']
        rows = len(table_data)
        cols = len(table_data[0]) if rows > 0 else 0
        
        if rows == 0 or cols == 0:
            logger.error("Table must have at least one row and one column")
            return None
        
        # Get position
        left = left or Inches(element_data.get('left', 1))
        top = top or Inches(element_data.get('top', 1))
        width = width or Inches(element_data.get('width', 8))
        height = height or Inches(element_data.get('height', rows * 0.5))
        
        # Create table
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Fill table data
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                if j < cols:  # Ensure we don't exceed the number of columns
                    cell = table.cell(i, j)
                    
                    if isinstance(cell_data, dict):
                        # Complex cell with formatting
                        cell.text = cell_data.get('text', '')
                        
                        # Apply cell style
                        if 'style' in cell_data:
                            style = cell_data['style']
                            
                            # Fill color
                            if 'fill_color' in style:
                                fill_color = get_rgb_color(style['fill_color'])
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(*fill_color)
                            
                            # Text formatting
                            format_text_frame(
                                cell.text_frame,
                                font=style.get('font', theme_settings['body_font']),
                                size=style.get('size', theme_settings['body_font_size']),
                                color=get_rgb_color(style.get('color', theme_settings['text_color'])),
                                alignment=style.get('align', PP_ALIGN.LEFT),
                                bold=style.get('bold', False),
                                italic=style.get('italic', False),
                                underline=style.get('underline', False)
                            )
                    else:
                        # Simple cell with just text
                        cell.text = str(cell_data)
                        
                        # Apply default formatting
                        format_text_frame(
                            cell.text_frame,
                            font=theme_settings['body_font'],
                            size=theme_settings['body_font_size'],
                            color=theme_settings['text_color']
                        )
                        
                        # Format header row differently
                        if i == 0 and element_data.get('has_header', True):
                            format_text_frame(
                                cell.text_frame,
                                bold=True,
                                color=theme_settings['accent_color']
                            )
        
        # Apply table styling
        if 'style' in element_data:
            style = element_data['style']
            
            # Border styles (TODO: python-pptx has limited support for this)
            # Zebra striping for rows
            if style.get('zebra_striping', False):
                for i in range(1, rows, 2):  # Start from second row (index 1)
                    for j in range(cols):
                        cell = table.cell(i, j)
                        cell.fill.solid()
                        
                        # Use lighter version of accent color or a specified alternate color
                        if 'alternate_row_color' in style:
                            alt_color = get_rgb_color(style['alternate_row_color'])
                        else:
                            # Lighten accent color
                            accent = theme_settings['accent_color']
                            alt_color = tuple(min(255, c + 40) for c in accent)
                        
                        cell.fill.fore_color.rgb = RGBColor(*alt_color)
        
        return table
    
    def create_chart(self, slide, left, top, width, height, element_data, theme_settings):
        """
        Create a chart element.
        
        Args:
            slide: The slide to add the chart to.
            left: Left position (can be None if specified in element_data).
            top: Top position (can be None if specified in element_data).
            width: Width (can be None).
            height: Height (can be None).
            element_data (dict): Chart configuration data.
            theme_settings (dict): Theme settings.
            
        Returns:
            The created chart shape.
        """
        # Get chart data
        if 'data' not in element_data or not element_data['data']:
            logger.error("Chart data not specified or empty")
            return None
        
        chart_data = element_data['data']
        
        # Get chart type
        chart_type_str = element_data.get('chart_type', 'bar').lower()
        chart_type_map = {
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'doughnut': XL_CHART_TYPE.DOUGHNUT,
            'area': XL_CHART_TYPE.AREA,
            'scatter': XL_CHART_TYPE.XY_SCATTER,
            'radar': XL_CHART_TYPE.RADAR
        }
        chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # Get position
        left = left or Inches(element_data.get('left', 1))
        top = top or Inches(element_data.get('top', 1))
        width = width or Inches(element_data.get('width', 8))
        height = height or Inches(element_data.get('height', 5))
        
        # Create chart data object
        chart_data_obj = CategoryChartData()
        
        # Add categories
        if 'categories' in chart_data:
            categories = chart_data['categories']
            chart_data_obj.categories = categories
        
        # Add series
        if 'series' in chart_data:
            for series in chart_data['series']:
                name = series.get('name', '')
                values = series.get('values', [])
                chart_data_obj.add_series(name, values)
        
        # Create chart
        chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data_obj).chart
        
        # Set chart title
        if 'title' in element_data:
            chart.has_title = True
            chart.chart_title.text_frame.text = element_data['title']
            
            # Format title
            format_text_frame(
                chart.chart_title.text_frame,
                font=theme_settings['body_font'],
                size=theme_settings['body_font_size'] + 2,  # Slightly larger than body
                bold=True,
                color=theme_settings['title_color']
            )
        
        # Set chart style properties
        if 'style' in element_data:
            style = element_data['style']
            
            # TODO: Apply additional chart styling when python-pptx
            # provides more support for chart customization
        
        return chart
    
    def create_code_block(self, slide, left, top, width, height, element_data, theme_settings):
        """
        Create a code block element (implemented as a formatted text box).
        
        Args:
            slide: The slide to add the code block to.
            left: Left position (can be None if specified in element_data).
            top: Top position (can be None if specified in element_data).
            width: Width (can be None).
            height: Height (can be None).
            element_data (dict): Code block configuration data.
            theme_settings (dict): Theme settings.
            
        Returns:
            The created text box shape.
        """
        # Get code content
        if 'code' not in element_data:
            logger.error("Code content not specified")
            return None
        
        code = element_data['code']
        
        # Get position
        left = left or Inches(element_data.get('left', 1))
        top = top or Inches(element_data.get('top', 1))
        width = width or Inches(element_data.get('width', 8))
        height = height or Inches(element_data.get('height', 4))
        
        # Create text box for code
        code_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = code_box.text_frame
        
        # Clear any existing text
        text_frame.clear()
        
        # Set code content
        p = text_frame.add_paragraph()
        p.text = code
        
        # Get code style properties
        font = element_data.get('font', 'Consolas')
        size = element_data.get('size', 12)
        
        # Default to light green color for code
        if 'color' in element_data:
            color = get_rgb_color(element_data['color'])
        else:
            color = (0, 200, 0)  # Light green default for code
        
        # Apply formatting
        format_text_frame(
            text_frame,
            font=font,
            size=size,
            color=color,
            line_spacing=1.2  # Slightly increased line spacing for code readability
        )
        
        # Add background rectangle for code block
        if element_data.get('background', True):
            # Create rectangle slightly larger than text box
            bg_left = left - Inches(0.1)
            bg_top = top - Inches(0.1)
            bg_width = width + Inches(0.2)
            bg_height = height + Inches(0.2)
            
            bg_rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 
                                           bg_left, bg_top, bg_width, bg_height)
            
            # Set rectangle style
            bg_color = get_rgb_color(element_data.get('background_color', '#202020'))
            bg_rect.fill.solid()
            bg_rect.fill.fore_color.rgb = RGBColor(*bg_color)
            
            # Add subtle outline
            bg_rect.line.color.rgb = RGBColor(100, 100, 100)
            bg_rect.line.width = Pt(1)
            
            # Send background to back
            slide.shapes._spTree.remove(bg_rect._element)
            slide.shapes._spTree.insert(0, bg_rect._element)
        
        return code_box