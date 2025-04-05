"""
PowerPoint Presentation Generator for Strings in C++ Presentation

This module contains functionality for generating a PowerPoint presentation about
Strings in C++ from a YAML file with detailed content.
"""

import os
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class StringsPPTGenerator:
    """
    A class for generating a PowerPoint presentation about Strings in C++.
    """
    
    def __init__(self, template_path=None):
        """
        Initialize the StringsPPTGenerator with an optional template.
        
        Args:
            template_path (str, optional): Path to a PowerPoint template file.
        """
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        # Default styling for dark theme
        self.default_style = {
            'title_font': 'Segoe UI Light',
            'title_size': 44,
            'subtitle_font': 'Segoe UI',
            'subtitle_size': 32,
            'body_font': 'Segoe UI',
            'body_size': 20,
            'code_font': 'Consolas',
            'code_size': 16,
            'background_color': RGBColor(26, 26, 26),  # Dark charcoal (#1A1A1A)
            'title_color': RGBColor(0, 255, 255),      # Cyan (#00FFFF)
            'subtitle_color': RGBColor(187, 134, 252), # Light Purple (#BB86FC)
            'text_color': RGBColor(224, 224, 224),     # Light Gray (#E0E0E0)
            'code_color': RGBColor(3, 218, 198),       # Light Green (#03DAC6)
            'accent_color': RGBColor(0, 191, 255),     # Light Blue (#00BFFF)
            'highlight_color': RGBColor(255, 140, 0),  # Orange (#FF8C00)
            'dark_gray': RGBColor(45, 45, 45),         # Dark Gray (#2D2D2D)
            'very_dark_gray': RGBColor(37, 37, 37)     # Very Dark Gray (#252525)
        }
    
    def generate_from_file(self, input_file_path, output_path):
        """
        Generate a PowerPoint presentation from a YAML file.
        
        Args:
            input_file_path (str): Path to the YAML input file.
            output_path (str): Path where the PowerPoint file should be saved.
            
        Returns:
            bool: True if successful, False otherwise.
        """
        try:
            # Read and parse the input file
            with open(input_file_path, 'r', encoding='utf-8') as f:
                content = yaml.safe_load(f)
            
            # Apply presentation-wide settings
            self._apply_presentation_settings(content.get('settings', {}))
            
            # Process variables
            if 'variables' in content:
                self._apply_variables(content.get('variables', {}))
            
            # Process slides
            for slide_data in content.get('slides', []):
                self._create_slide(slide_data)
            
            # Save the presentation
            self.prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"Error generating presentation: {e}")
            return False
    
    def _apply_presentation_settings(self, settings):
        """
        Apply presentation-wide settings.
        
        Args:
            settings (dict): Dictionary of presentation settings.
        """
        # Apply custom styling if provided
        if 'style' in settings:
            style = settings['style']
            if 'title_font' in style:
                self.default_style['title_font'] = style['title_font']
            if 'title_size' in style:
                self.default_style['title_size'] = style['title_size']
            if 'body_font' in style:
                self.default_style['body_font'] = style['body_font']
            if 'body_size' in style:
                self.default_style['body_size'] = style['body_size']
            if 'theme_color' in style:
                color = style['theme_color']
                if isinstance(color, list) and len(color) == 3:
                    self.default_style['accent_color'] = RGBColor(*color)
    
    def _apply_variables(self, variables):
        """
        Apply variables from the input file.
        
        Args:
            variables (dict): Dictionary of variables.
        """
        # Apply color variables if provided
        for color_name in ['background_color', 'title_color', 'subtitle_color', 
                          'text_color', 'code_color', 'accent_color', 
                          'highlight_color', 'dark_gray', 'very_dark_gray']:
            if color_name in variables:
                color = variables[color_name]
                if isinstance(color, list) and len(color) == 3:
                    self.default_style[color_name] = RGBColor(*color)
        
        # Store all variables for text replacement
        self.variables = variables
    
    def _create_slide(self, slide_data):
        """
        Create a slide based on the provided data.
        
        Args:
            slide_data (dict): Dictionary containing slide data.
        """
        slide_type = slide_data.get('type', 'title_and_content')
        
        # Select the appropriate slide layout
        layout = self._get_slide_layout(slide_type)
        slide = self.prs.slides.add_slide(layout)
        
        # Apply background color to the slide
        self._apply_background_color(slide)
        
        # Process slide content based on type
        if slide_type == 'title_slide':
            self._create_title_slide(slide, slide_data)
        elif slide_type == 'title_and_content':
            self._create_title_content_slide(slide, slide_data)
        elif slide_type == 'section':
            self._create_section_slide(slide, slide_data)
        elif slide_type == 'title_only':
            self._create_title_only_slide(slide, slide_data)
        elif slide_type == 'blank':
            self._create_blank_slide(slide, slide_data)
        else:
            # Default to title and content
            self._create_title_content_slide(slide, slide_data)
    
    def _apply_background_color(self, slide):
        """
        Apply background color to a slide.
        
        Args:
            slide (Slide): The slide to modify.
        """
        # Get the slide background
        background = slide.background
        
        # Apply fill
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.default_style['background_color']
    
    def _get_slide_layout(self, slide_type):
        """
        Get the appropriate slide layout for the given slide type.
        
        Args:
            slide_type (str): Type of slide to create.
            
        Returns:
            SlideLayout: The slide layout to use.
        """
        # Map slide types to indices in the slide layout collection
        slide_layouts = {
            'title_slide': 0,       # Title Slide
            'title_and_content': 1, # Title and Content
            'section': 2,           # Section Header
            'two_content': 3,       # Two Content
            'title_only': 5,        # Title Only
            'blank': 6              # Blank
        }
        
        layout_idx = slide_layouts.get(slide_type, 1)  # Default to title and content
        return self.prs.slide_layouts[layout_idx]
    
    def _create_title_slide(self, slide, slide_data):
        """
        Create a title slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['title_color'])
        
        if 'subtitle' in slide_data and hasattr(slide.placeholders, '__iter__'):
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # Subtitle placeholder
                    shape.text = self._replace_variables(slide_data['subtitle'])
                    self._apply_text_formatting(shape.text_frame,
                                              font=self.default_style['subtitle_font'],
                                              size=self.default_style['subtitle_size'],
                                              color=self.default_style['subtitle_color'])
                    break
    
    def _create_title_content_slide(self, slide, slide_data):
        """
        Create a title and content slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            
            # Special formatting for section titles vs regular titles
            if slide_data.get('type') == 'section':
                self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['accent_color'])
            else:
                self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['title_color'])
        
        if 'content' in slide_data:
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Content placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                if isinstance(slide_data['content'], list):
                    # Create a bulleted list
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for i, item in enumerate(slide_data['content']):
                        p = text_frame.add_paragraph()
                        p.text = self._replace_variables(item)
                        p.level = 0  # Top level bullet
                        self._apply_text_formatting(p, 
                                                  font=self.default_style['body_font'],
                                                  size=self.default_style['body_size'],
                                                  color=self.default_style['text_color'])
                else:
                    # Process content as text with possible code blocks
                    content = self._replace_variables(slide_data['content'])
                    
                    # Check for code blocks surrounded by ```
                    if "```" in content:
                        self._process_content_with_code(content_placeholder, content)
                    else:
                        # Regular text content
                        content_placeholder.text = content
                        self._apply_text_formatting(content_placeholder.text_frame,
                                                  font=self.default_style['body_font'],
                                                  size=self.default_style['body_size'],
                                                  color=self.default_style['text_color'])
    
    def _process_content_with_code(self, placeholder, content):
        """
        Process content that contains code blocks marked with triple backticks.
        
        Args:
            placeholder: The placeholder shape to populate.
            content (str): The content text with code blocks.
        """
        # Clear existing content
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        # Split by code blocks
        parts = content.split("```")
        
        # Process each part
        for i, part in enumerate(parts):
            # Skip empty parts
            if not part.strip():
                continue
                
            if i % 2 == 1:  # This is a code block
                # Extract language if specified on first line
                lines = part.strip().split('\n')
                if len(lines) > 1 and not lines[0].strip().startswith('#'):
                    code = '\n'.join(lines[1:])  # Skip language line
                else:
                    code = part
                
                # Add code paragraph
                p = text_frame.add_paragraph()
                p.text = code
                self._apply_text_formatting(p, 
                                          font=self.default_style['code_font'],
                                          size=self.default_style['code_size'],
                                          color=self.default_style['code_color'])
            else:  # This is regular text
                # Split into lines and add each as a paragraph
                lines = part.strip().split('\n')
                for line in lines:
                    if line.strip():
                        p = text_frame.add_paragraph()
                        p.text = line
                        self._apply_text_formatting(p, 
                                                  font=self.default_style['body_font'],
                                                  size=self.default_style['body_size'],
                                                  color=self.default_style['text_color'])
    
    def _create_section_slide(self, slide, slide_data):
        """
        Create a section header slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['accent_color'])
    
    def _create_title_only_slide(self, slide, slide_data):
        """
        Create a title only slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['title_color'])
        
        # Process any custom elements
        if 'elements' in slide_data:
            for element in slide_data['elements']:
                self._add_custom_element(slide, element)
    
    def _create_blank_slide(self, slide, slide_data):
        """
        Create a blank slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        # Process any custom elements
        if 'elements' in slide_data:
            for element in slide_data['elements']:
                self._add_custom_element(slide, element)
    
    def _add_custom_element(self, slide, element):
        """
        Add a custom element to the slide.
        
        Args:
            slide (Slide): The slide to add the element to.
            element (dict): Element data.
        """
        element_type = element.get('type', 'text_box')
        
        if element_type == 'text_box':
            left = Inches(element.get('left', 1))
            top = Inches(element.get('top', 1))
            width = Inches(element.get('width', 4))
            height = Inches(element.get('height', 1))
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            
            # Set text content
            text_frame = textbox.text_frame
            text_frame.clear()
            
            # Check if the text contains a table
            text = self._replace_variables(element.get('text', ''))
            if '|' in text and '-|-' in text:  # Simple table detection
                self._add_markdown_table(text_frame, text, 
                                       element.get('font', self.default_style['code_font']),
                                       element.get('size', self.default_style['code_size']))
            else:
                # Regular text
                textbox.text = text
                self._apply_text_formatting(text_frame, 
                                         font=element.get('font', self.default_style['body_font']),
                                         size=element.get('size', self.default_style['body_size']),
                                         color=self.default_style['text_color'])
    
    def _add_markdown_table(self, text_frame, markdown_text, font, size):
        """
        Add a markdown-style table to a text frame.
        
        Args:
            text_frame: The text frame to add the table to.
            markdown_text (str): The markdown table text.
            font (str): The font to use.
            size (int): The font size to use.
        """
        # Split the markdown text into lines
        lines = markdown_text.strip().split('\n')
        
        # Process each line
        for i, line in enumerate(lines):
            # Skip separator lines (like |-|-|)
            if line.strip().replace('-', '').replace('|', '').strip() == '':
                continue
                
            # Clean the line and split into cells
            cleaned_line = line.strip()
            if cleaned_line.startswith('|'):
                cleaned_line = cleaned_line[1:]
            if cleaned_line.endswith('|'):
                cleaned_line = cleaned_line[:-1]
                
            cells = [cell.strip() for cell in cleaned_line.split('|')]
            
            # Create a paragraph for this row
            p = text_frame.add_paragraph()
            
            # Format header rows differently
            if i == 0:  # Header row
                p.text = ' | '.join(cells)
                self._apply_text_formatting(p, 
                                         font=font,
                                         size=size,
                                         color=self.default_style['accent_color'],
                                         bold=True)
            else:  # Data rows
                # Alternate row colors for readability
                p.text = ' | '.join(cells)
                if i % 2 == 0:  # Even rows
                    self._apply_text_formatting(p, 
                                             font=font,
                                             size=size,
                                             color=self.default_style['text_color'])
                else:  # Odd rows
                    self._apply_text_formatting(p, 
                                             font=font,
                                             size=size,
                                             color=self.default_style['text_color'])
    
    def _apply_text_formatting(self, text_frame, font=None, size=None, color=None, alignment=None, bold=False):
        """
        Apply formatting to text.
        
        Args:
            text_frame: The text frame to format.
            font (str, optional): Font name.
            size (int, optional): Font size in points.
            color (RGBColor, optional): Font color.
            alignment (PP_ALIGN, optional): Text alignment.
            bold (bool, optional): Whether to make the text bold.
        """
        if not hasattr(text_frame, 'paragraphs'):
            return
            
        for paragraph in text_frame.paragraphs:
            if alignment is not None:
                paragraph.alignment = alignment
            
            for run in paragraph.runs:
                if font is not None:
                    run.font.name = font
                
                if size is not None:
                    run.font.size = Pt(size)
                
                if color is not None:
                    run.font.color.rgb = color
                    
                if bold:
                    run.font.bold = True
    
    def _replace_variables(self, text):
        """
        Replace variables in text with their values.
        
        Args:
            text (str): Text that may contain variables.
            
        Returns:
            str: Text with variables replaced.
        """
        if not isinstance(text, str):
            return text
            
        result = text
        if hasattr(self, 'variables'):
            for var_name, var_value in self.variables.items():
                placeholder = '{{' + var_name + '}}'
                if isinstance(var_value, (str, int, float)):
                    result = result.replace(placeholder, str(var_value))
        
        return result

def main():
    """
    Example usage of the StringsPPTGenerator class.
    """
    # Create a PPT generator
    generator = StringsPPTGenerator()
    
    # Generate a presentation from a YAML file
    success = generator.generate_from_file('strings_presentation.yaml', 'strings_presentation.pptx')
    
    if success:
        print("Strings in C++ presentation generated successfully!")
    else:
        print("Failed to generate presentation.")

if __name__ == "__main__":
    main()