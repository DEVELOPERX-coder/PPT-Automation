"""
PowerPoint Presentation Generator with Dark Theme Support

This module contains functionality for generating PowerPoint presentations
from custom-formatted input files with dark theme support.
"""

import os
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PPTGenerator:
    """
    A class for generating PowerPoint presentations from custom input files.
    """
    
    def __init__(self, template_path=None, dark_mode=True):
        """
        Initialize the PPTGenerator with an optional template and theme option.
        
        Args:
            template_path (str, optional): Path to a PowerPoint template file.
            dark_mode (bool, optional): Whether to use dark mode. Defaults to True.
        """
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.dark_mode = dark_mode
        
        # Default styling
        if self.dark_mode:
            # Dark theme defaults
            self.default_style = {
                'title_font': 'Arial',
                'title_size': 44,
                'body_font': 'Arial',
                'body_size': 18,
                'theme_color': RGBColor(0, 191, 255),  # Deep sky blue
                'background_color': RGBColor(30, 30, 30),  # Dark gray
                'text_color': RGBColor(220, 220, 220),  # Light gray
                'accent_color': RGBColor(0, 191, 255),  # Deep sky blue
                'highlight_color': RGBColor(255, 140, 0)  # Orange
            }
        else:
            # Light theme defaults
            self.default_style = {
                'title_font': 'Calibri',
                'title_size': 44,
                'body_font': 'Calibri',
                'body_size': 18,
                'theme_color': RGBColor(0, 112, 192),  # Blue
                'background_color': RGBColor(255, 255, 255),  # White
                'text_color': RGBColor(0, 0, 0),  # Black
                'accent_color': RGBColor(0, 112, 192),  # Blue
                'highlight_color': RGBColor(255, 140, 0)  # Orange
            }
    
    def generate_from_file(self, input_file_path, output_path):
        """
        Generate a PowerPoint presentation from a custom input file.
        
        Args:
            input_file_path (str): Path to the input file.
            output_path (str): Path where the PowerPoint file should be saved.
            
        Returns:
            bool: True if successful, False otherwise.
        """
        try:
            # Read and parse the input file
            with open(input_file_path, 'r', encoding='utf-8') as f:
                content = yaml.safe_load(f)
            
            # Apply presentation-wide settings
            self._apply_presentation_settings(content.get('settings', {}))
            
            # Process variables
            if 'variables' in content:
                self._apply_variables(content.get('variables', {}))
            
            # Process slides
            for slide_data in content.get('slides', []):
                self._create_slide(slide_data)
            
            # Save the presentation
            self.prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"Error generating presentation: {e}")
            return False
    
    def _apply_presentation_settings(self, settings):
        """
        Apply presentation-wide settings.
        
        Args:
            settings (dict): Dictionary of presentation settings.
        """
        # Apply custom styling if provided
        if 'style' in settings:
            style = settings['style']
            if 'title_font' in style:
                self.default_style['title_font'] = style['title_font']
            if 'title_size' in style:
                self.default_style['title_size'] = style['title_size']
            if 'body_font' in style:
                self.default_style['body_font'] = style['body_font']
            if 'body_size' in style:
                self.default_style['body_size'] = style['body_size']
            if 'theme_color' in style:
                color = style['theme_color']
                if isinstance(color, list) and len(color) == 3:
                    self.default_style['theme_color'] = RGBColor(*color)
                    self.default_style['accent_color'] = RGBColor(*color)
    
    def _apply_variables(self, variables):
        """
        Apply variables from the input file.
        
        Args:
            variables (dict): Dictionary of variables.
        """
        # Apply color variables if provided
        if 'background_color' in variables:
            color = variables['background_color']
            if isinstance(color, list) and len(color) == 3:
                self.default_style['background_color'] = RGBColor(*color)
        
        if 'text_color' in variables:
            color = variables['text_color']
            if isinstance(color, list) and len(color) == 3:
                self.default_style['text_color'] = RGBColor(*color)
        
        if 'accent_color' in variables:
            color = variables['accent_color']
            if isinstance(color, list) and len(color) == 3:
                self.default_style['accent_color'] = RGBColor(*color)
        
        if 'highlight_color' in variables:
            color = variables['highlight_color']
            if isinstance(color, list) and len(color) == 3:
                self.default_style['highlight_color'] = RGBColor(*color)
        
        # Store all variables for text replacement
        self.variables = variables
    
    def _create_slide(self, slide_data):
        """
        Create a slide based on the provided data.
        
        Args:
            slide_data (dict): Dictionary containing slide data.
        """
        slide_type = slide_data.get('type', 'title_and_content')
        
        # Select the appropriate slide layout
        layout = self._get_slide_layout(slide_type)
        slide = self.prs.slides.add_slide(layout)
        
        # Apply background color to the slide
        self._apply_background_color(slide)
        
        # Process slide content based on type
        if slide_type == 'title_slide':
            self._create_title_slide(slide, slide_data)
        elif slide_type == 'title_and_content':
            self._create_title_content_slide(slide, slide_data)
        elif slide_type == 'section':
            self._create_section_slide(slide, slide_data)
        elif slide_type == 'two_content':
            self._create_two_content_slide(slide, slide_data)
        elif slide_type == 'title_only':
            self._create_title_only_slide(slide, slide_data)
        elif slide_type == 'blank':
            self._create_blank_slide(slide, slide_data)
        else:
            # Default to title and content
            self._create_title_content_slide(slide, slide_data)
    
    def _apply_background_color(self, slide):
        """
        Apply background color to a slide.
        
        Args:
            slide (Slide): The slide to modify.
        """
        # Get the slide background
        background = slide.background
        
        # Apply fill
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.default_style['background_color']
    
    def _get_slide_layout(self, slide_type):
        """
        Get the appropriate slide layout for the given slide type.
        
        Args:
            slide_type (str): Type of slide to create.
            
        Returns:
            SlideLayout: The slide layout to use.
        """
        # Map slide types to indices in the slide layout collection
        slide_layouts = {
            'title_slide': 0,       # Title Slide
            'title_and_content': 1, # Title and Content
            'section': 2,           # Section Header
            'two_content': 3,       # Two Content
            'title_only': 5,        # Title Only
            'blank': 6              # Blank
        }
        
        layout_idx = slide_layouts.get(slide_type, 1)  # Default to title and content
        return self.prs.slide_layouts[layout_idx]
    
    def _create_title_slide(self, slide, slide_data):
        """
        Create a title slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['text_color'])
        
        if 'subtitle' in slide_data and hasattr(slide.placeholders, '__iter__'):
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # Subtitle placeholder
                    shape.text = self._replace_variables(slide_data['subtitle'])
                    self._apply_text_formatting(shape.text_frame,
                                              font=self.default_style['body_font'],
                                              size=self.default_style['body_size'],
                                              color=self.default_style['text_color'])
                    break
    
    def _create_title_content_slide(self, slide, slide_data):
        """
        Create a title and content slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['text_color'])
        
        if 'content' in slide_data:
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Content placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                if isinstance(slide_data['content'], list):
                    # Create a bulleted list
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for i, item in enumerate(slide_data['content']):
                        p = text_frame.add_paragraph()
                        p.text = self._replace_variables(item)
                        p.level = 0  # Top level bullet
                        self._apply_text_formatting(p, 
                                                  font=self.default_style['body_font'],
                                                  size=self.default_style['body_size'],
                                                  color=self.default_style['text_color'])
                else:
                    # Simple text content or code block
                    content = self._replace_variables(slide_data['content'])
                    
                    # Check if it's a code block
                    if content.strip().startswith('```') and content.strip().endswith('```'):
                        # Extract code and format it
                        lines = content.strip().split('\n')
                        code_lang = lines[0].replace('```', '').strip()
                        code = '\n'.join(lines[1:-1])
                        
                        # Add code as fixed-width text
                        text_frame = content_placeholder.text_frame
                        text_frame.clear()
                        p = text_frame.add_paragraph()
                        p.text = code
                        self._apply_text_formatting(p, 
                                                  font='Consolas',  # Use monospace font for code
                                                  size=14,  # Smaller size for code
                                                  color=self.default_style['text_color'])
                    else:
                        # Regular text
                        content_placeholder.text = content
                        self._apply_text_formatting(content_placeholder.text_frame,
                                                  font=self.default_style['body_font'],
                                                  size=self.default_style['body_size'],
                                                  color=self.default_style['text_color'])
        
        # Process any images
        if 'image' in slide_data:
            self._add_image_to_slide(slide, slide_data['image'])
    
    def _create_section_slide(self, slide, slide_data):
        """
        Create a section header slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['accent_color'])  # Use accent color for section headers
    
    def _create_two_content_slide(self, slide, slide_data):
        """
        Create a two content slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['text_color'])
        
        # Find the content placeholders
        left_placeholder = None
        right_placeholder = None
        
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 7:  # Content placeholder
                if left_placeholder is None:
                    left_placeholder = shape
                else:
                    right_placeholder = shape
                    break
        
        # Add left content
        if left_placeholder and 'left_content' in slide_data:
            if isinstance(slide_data['left_content'], list):
                # Create a bulleted list
                text_frame = left_placeholder.text_frame
                text_frame.clear()
                
                for item in slide_data['left_content']:
                    p = text_frame.add_paragraph()
                    p.text = self._replace_variables(item)
                    p.level = 0  # Top level bullet
                    self._apply_text_formatting(p, 
                                              font=self.default_style['body_font'],
                                              size=self.default_style['body_size'],
                                              color=self.default_style['text_color'])
            else:
                # Simple text content
                left_placeholder.text = self._replace_variables(slide_data['left_content'])
                self._apply_text_formatting(left_placeholder.text_frame,
                                          font=self.default_style['body_font'],
                                          size=self.default_style['body_size'],
                                          color=self.default_style['text_color'])
        
        # Add right content
        if right_placeholder and 'right_content' in slide_data:
            if isinstance(slide_data['right_content'], list):
                # Create a bulleted list
                text_frame = right_placeholder.text_frame
                text_frame.clear()
                
                for item in slide_data['right_content']:
                    p = text_frame.add_paragraph()
                    p.text = self._replace_variables(item)
                    p.level = 0  # Top level bullet
                    self._apply_text_formatting(p, 
                                              font=self.default_style['body_font'],
                                              size=self.default_style['body_size'],
                                              color=self.default_style['text_color'])
            else:
                # Simple text content
                right_placeholder.text = self._replace_variables(slide_data['right_content'])
                self._apply_text_formatting(right_placeholder.text_frame,
                                          font=self.default_style['body_font'],
                                          size=self.default_style['body_size'],
                                          color=self.default_style['text_color'])
        
        # Process any images
        if 'left_image' in slide_data and left_placeholder:
            self._add_image_to_placeholder(left_placeholder, slide_data['left_image'])
            
        if 'right_image' in slide_data and right_placeholder:
            self._add_image_to_placeholder(right_placeholder, slide_data['right_image'])
    
    def _create_title_only_slide(self, slide, slide_data):
        """
        Create a title only slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = self._replace_variables(slide_data['title'])
            self._apply_text_formatting(title.text_frame, 
                                       font=self.default_style['title_font'],
                                       size=self.default_style['title_size'],
                                       color=self.default_style['text_color'])
        
        # Process any custom elements
        if 'elements' in slide_data:
            for element in slide_data['elements']:
                self._add_custom_element(slide, element)
    
    def _create_blank_slide(self, slide, slide_data):
        """
        Create a blank slide with the given data.
        
        Args:
            slide (Slide): The slide to populate.
            slide_data (dict): Data for the slide.
        """
        # Process any custom elements
        if 'elements' in slide_data:
            for element in slide_data['elements']:
                self._add_custom_element(slide, element)
    
    def _add_custom_element(self, slide, element):
        """
        Add a custom element to the slide.
        
        Args:
            slide (Slide): The slide to add the element to.
            element (dict): Element data.
        """
        element_type = element.get('type', 'text_box')
        
        if element_type == 'text_box':
            left = Inches(element.get('left', 1))
            top = Inches(element.get('top', 1))
            width = Inches(element.get('width', 4))
            height = Inches(element.get('height', 1))
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text = self._replace_variables(element.get('text', ''))
            
            font = element.get('font', self.default_style['body_font'])
            size = element.get('size', self.default_style['body_size'])
            self._apply_text_formatting(textbox.text_frame, 
                                       font=font, 
                                       size=size,
                                       color=self.default_style['text_color'])
        
        elif element_type == 'image':
            self._add_image_to_slide(slide, element)
    
    def _add_image_to_slide(self, slide, image_data):
        """
        Add an image to the slide.
        
        Args:
            slide (Slide): The slide to add the image to.
            image_data (dict): Image data.
        """
        image_path = image_data.get('path', '')
        
        if not os.path.exists(image_path):
            print(f"Image not found: {image_path}")
            return
        
        left = Inches(image_data.get('left', 1))
        top = Inches(image_data.get('top', 1))
        width = Inches(image_data.get('width', 4))
        height = Inches(image_data.get('height', 3))
        
        slide.shapes.add_picture(image_path, left, top, width, height)
    
    def _add_image_to_placeholder(self, placeholder, image_data):
        """
        Add an image to a placeholder.
        
        Args:
            placeholder (Shape): The placeholder to add the image to.
            image_data (dict): Image data.
        """
        image_path = image_data.get('path', '')
        
        if not os.path.exists(image_path):
            print(f"Image not found: {image_path}")
            return
        
        # Clear any existing content
        if hasattr(placeholder, 'text'):
            placeholder.text = ''
        
        # Add image to the placeholder
        placeholder.insert_picture(image_path)
    
    def _apply_text_formatting(self, text_frame, font=None, size=None, color=None, alignment=None):
        """
        Apply formatting to text.
        
        Args:
            text_frame: The text frame to format.
            font (str, optional): Font name.
            size (int, optional): Font size in points.
            color (RGBColor, optional): Font color.
            alignment (PP_ALIGN, optional): Text alignment.
        """
        if not hasattr(text_frame, 'paragraphs'):
            return
            
        for paragraph in text_frame.paragraphs:
            if alignment is not None:
                paragraph.alignment = alignment
            
            for run in paragraph.runs:
                if font is not None:
                    run.font.name = font
                
                if size is not None:
                    run.font.size = Pt(size)
                
                if color is not None:
                    run.font.color.rgb = color
    
    def _replace_variables(self, text):
        """
        Replace variables in text with their values.
        
        Args:
            text (str): Text that may contain variables.
            
        Returns:
            str: Text with variables replaced.
        """
        if not isinstance(text, str):
            return text
            
        result = text
        if hasattr(self, 'variables'):
            for var_name, var_value in self.variables.items():
                placeholder = '{{' + var_name + '}}'
                if isinstance(var_value, (str, int, float)):
                    result = result.replace(placeholder, str(var_value))
        
        return result

def main():
    """
    Example usage of the PPTGenerator class.
    """
    # Create a PPT generator with dark mode
    generator = PPTGenerator(dark_mode=True)
    
    # Generate a presentation from a YAML file
    success = generator.generate_from_file('dsa_day5_presentation.yaml', 'dsa_day5_presentation.pptx')
    
    if success:
        print("Presentation generated successfully!")
    else:
        print("Failed to generate presentation.")

if __name__ == "__main__":
    main()