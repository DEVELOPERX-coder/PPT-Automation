#!/usr/bin/env python3
# md2pptx.py - Enhanced Markdown to PowerPoint Converter
import os
import re
import sys
import math
import argparse
import traceback
import yaml
from pathlib import Path
from enum import Enum

class BackendType(Enum):
    """Enum defining the presentation backend types."""
    PYTHON_PPTX = "python-pptx"
    WIN32COM = "win32com"

def parse_markdown(md_file):
    """Parse the markdown file into slide data.
    
    Args:
        md_file: Path to the markdown file
        
    Returns:
        tuple: (slides, settings) where slides is a list of slide data dictionaries
               and settings is a dictionary of global settings
    """
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    settings = {}
    
    # Extract global settings if present (YAML frontmatter)
    if content.startswith('---'):
        end_index = content.find('---', 3)
        if end_index != -1:
            yaml_content = content[3:end_index].strip()
            try:
                settings = yaml.safe_load(yaml_content)
                print(f"Loaded global settings: {len(settings)} properties")
            except Exception as e:
                print(f"Warning: Could not parse global settings: {e}")
            content = content[end_index+3:].strip()
    
    # Split into slides by headings
    slide_blocks = re.split(r'\n\s*#\s+', content)
    
    # If content doesn't start with a heading, process differently
    if not content.lstrip().startswith('#'):
        slide_blocks = slide_blocks[1:]  # Skip non-slide content before first heading
    
    # Process each slide
    for block in slide_blocks:
        if not block.strip():
            continue
            
        # Get slide title
        title_match = re.match(r'([^\n]+)', block)
        title = title_match.group(1).strip() if title_match else ""
        block = block[title_match.end():].strip() if title_match else block.strip()
        
        # Get slide settings
        slide_settings = {}
        settings_match = re.match(r'\{\s*(.*?)\s*\}', block, re.DOTALL)
        if settings_match:
            settings_str = settings_match.group(1)
            # Parse as YAML first (more robust)
            try:
                yaml_str = "{\n" + settings_str + "\n}"
                slide_settings = yaml.safe_load(yaml_str)
            except:
                # Fallback to regex parsing for key-value pairs
                for setting in re.finditer(r'(\w+)\s*:\s*("[^"]*"|\'[^\']*\'|[^,}\s][^,}]*)', settings_str):
                    key = setting.group(1).strip()
                    value = setting.group(2).strip()
                    # Strip quotes if present
                    if (value.startswith('"') and value.endswith('"')) or \
                       (value.startswith("'") and value.endswith("'")):
                        value = value[1:-1]
                    slide_settings[key] = value
                
            block = block[settings_match.end():].strip()
        
        # Extract elements
        elements = []
        for element_match in re.finditer(r':::(\w+)(?:\[(.*?)\])?\s*\n((?:\{.*?\})?\s*[\s\S]*?)(?=\n:::|$)', block):
            element_type = element_match.group(1)
            inline_props = element_match.group(2) or ""
            element_content = element_match.group(3).strip()
            
            # Parse properties
            properties = {}
            
            # From inline props
            if inline_props:
                try:
                    # Try YAML parsing first
                    yaml_str = "{\n" + inline_props + "\n}"
                    inline_props_dict = yaml.safe_load(yaml_str)
                    if inline_props_dict:
                        properties.update(inline_props_dict)
                except:
                    # Fallback to regex parsing
                    for prop in re.finditer(r'(\w+)\s*:\s*("[^"]*"|\'[^\']*\'|[^,\s][^,]*)', inline_props):
                        key, value = prop.group(1).strip(), prop.group(2).strip()
                        if (value.startswith('"') and value.endswith('"')) or \
                           (value.startswith("'") and value.endswith("'")):
                            value = value[1:-1]
                        properties[key] = value
            
            # From property block
            prop_block_match = re.match(r'\{\s*(.*?)\s*\}([\s\S]*)', element_content, re.DOTALL)
            if prop_block_match:
                prop_str = prop_block_match.group(1)
                try:
                    # Try YAML parsing first
                    yaml_str = "{\n" + prop_str + "\n}"
                    prop_dict = yaml.safe_load(yaml_str)
                    if prop_dict:
                        properties.update(prop_dict)
                except:
                    # Fallback to regex parsing
                    for prop in re.finditer(r'(\w+)\s*:\s*("[^"]*"|\'[^\']*\'|[^,}\s][^,}]*)', prop_str):
                        key, value = prop.group(1).strip(), prop.group(2).strip()
                        if (value.startswith('"') and value.endswith('"')) or \
                           (value.startswith("'") and value.endswith("'")):
                            value = value[1:-1]
                        properties[key] = value
                        
                element_content = prop_block_match.group(2).strip()
            
            elements.append({
                'type': element_type, 
                'properties': properties,
                'content': element_content
            })
        
        slides.append({
            'title': title,
            'settings': slide_settings,
            'elements': elements
        })
    
    print(f"Parsed {len(slides)} slides with a total of {sum(len(slide['elements']) for slide in slides)} elements")
    return slides, settings

def get_position_and_size(properties):
    """Extract position and size from properties with unit conversion.
    
    Args:
        properties: Dictionary of element properties
        
    Returns:
        tuple: (left, top, width, height) in points
    """
    # Default values (1 inch from top-left, 4x1 inches)
    left = 72    # 1 inch = 72 points
    top = 72     # 1 inch
    width = 288  # 4 inches
    height = 72  # 1 inch
    
    # Parse x/left position
    if 'x' in properties:
        x_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(properties['x']))
        if x_match:
            x_num = float(x_match.group(1))
            x_unit = x_match.group(3) or 'in'
            if x_unit == 'in': left = x_num * 72
            elif x_unit == 'cm': left = x_num * 28.35
            elif x_unit == 'pt': left = x_num
            elif x_unit == 'px': left = x_num * 0.75
    
    # Parse y/top position
    if 'y' in properties:
        y_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(properties['y']))
        if y_match:
            y_num = float(y_match.group(1))
            y_unit = y_match.group(3) or 'in'
            if y_unit == 'in': top = y_num * 72
            elif y_unit == 'cm': top = y_num * 28.35
            elif y_unit == 'pt': top = y_num
            elif y_unit == 'px': top = y_num * 0.75
    
    # Parse width
    if 'width' in properties:
        w_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(properties['width']))
        if w_match:
            w_num = float(w_match.group(1))
            w_unit = w_match.group(3) or 'in'
            if w_unit == 'in': width = w_num * 72
            elif w_unit == 'cm': width = w_num * 28.35
            elif w_unit == 'pt': width = w_num
            elif w_unit == 'px': width = w_num * 0.75
    
    # Parse height
    if 'height' in properties:
        h_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(properties['height']))
        if h_match:
            h_num = float(h_match.group(1))
            h_unit = h_match.group(3) or 'in'
            if h_unit == 'in': height = h_num * 72
            elif h_unit == 'cm': height = h_num * 28.35
            elif h_unit == 'pt': height = h_num
            elif h_unit == 'px': height = h_num * 0.75
    
    return left, top, width, height

def hex_to_rgb(hex_color):
    """Convert hex color to RGB integer.
    
    Args:
        hex_color: Hex color string like "#RRGGBB" or "#RGB"
        
    Returns:
        int: RGB color as integer
    """
    hex_color = hex_color.lstrip('#')
    
    # Convert shorthand #RGB to #RRGGBB
    if len(hex_color) == 3:
        hex_color = ''.join(c+c for c in hex_color)
    
    if len(hex_color) == 6:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        # Win32com expects BGR format
        return r + (g << 8) + (b << 16)
    
    return 0  # Black default

# Add this function after hex_to_rgb()
def debug_element(element_type, properties, content):
    """Print detailed debug information about an element to help troubleshoot issues."""
    print(f"\n  Adding {element_type} element:")
    print(f"  - Properties: {properties}")
    if content and len(content) > 50:
        print(f"  - Content: {content[:50]}...")
    else:
        print(f"  - Content: {content}")

def add_text_element_win32(slide, content, properties):
    """Add a text element to the slide using win32com.
    
    Args:
        slide: The PowerPoint slide object
        content: The text content
        properties: Dictionary of text properties
    """
    # Get position and size
    left, top, width, height = get_position_and_size(properties)
    
    # Create text box
    shape = slide.Shapes.AddTextbox(1, left, top, width, height)
    shape.TextFrame.TextRange.Text = content
    
    # Word wrap
    shape.TextFrame.WordWrap = True
    
    # Apply text formatting
    if 'font' in properties:
        shape.TextFrame.TextRange.Font.Name = properties['font']
    
    if 'font_size' in properties:
        try:
            size = float(properties['font_size'])
            shape.TextFrame.TextRange.Font.Size = size
        except ValueError:
            print(f"Warning: Invalid font size '{properties['font_size']}'")
    
    if 'font_color' in properties:
        color = properties['font_color']
        if color.startswith('#'):
            shape.TextFrame.TextRange.Font.Color.RGB = hex_to_rgb(color)
    
    if 'bold' in properties and properties['bold'].lower() in ('true', 'yes', '1'):
        shape.TextFrame.TextRange.Font.Bold = True
    
    if 'italic' in properties and properties['italic'].lower() in ('true', 'yes', '1'):
        shape.TextFrame.TextRange.Font.Italic = True
        
    if 'underline' in properties and properties['underline'].lower() in ('true', 'yes', '1'):
        shape.TextFrame.TextRange.Font.Underline = True
    
    if 'align' in properties:
        align = properties['align'].lower()
        if align == 'center':
            shape.TextFrame.TextRange.ParagraphFormat.Alignment = 2  # Center
        elif align == 'right':
            shape.TextFrame.TextRange.ParagraphFormat.Alignment = 3  # Right
        elif align == 'justify':
            shape.TextFrame.TextRange.ParagraphFormat.Alignment = 4  # Justify
    
    if 'vertical_align' in properties:
        valign = properties['vertical_align'].lower()
        if valign == 'middle':
            shape.TextFrame.VerticalAnchor = 3  # Middle
        elif valign == 'bottom':
            shape.TextFrame.VerticalAnchor = 2  # Bottom
    
    # Bullet points
    if 'bullet' in properties and properties['bullet'].lower() in ('true', 'yes', '1'):
        # Split content by lines and apply bullet to each paragraph
        lines = content.strip().split('\n')
        if len(lines) > 0:
            # Clear the textbox first
            shape.TextFrame.TextRange.Text = ""
            
            for i, line in enumerate(lines):
                if i > 0:
                    # Add a new paragraph after the first one
                    shape.TextFrame.TextRange.InsertAfter(line)
                    para = shape.TextFrame.TextRange.Paragraphs(i+1)
                else:
                    # Use the first paragraph
                    shape.TextFrame.TextRange.Text = line
                    para = shape.TextFrame.TextRange.Paragraphs(1)
                
                # Apply bullet formatting
                para.ParagraphFormat.Bullet.Type = 1  # Bullet
                
                # Apply bullet style if specified
                if 'bullet_style' in properties:
                    style = properties['bullet_style'].lower()
                    if style == 'number':
                        para.ParagraphFormat.Bullet.Type = 2  # Numbered
                    elif style == 'custom' and 'bullet_character' in properties:
                        para.ParagraphFormat.Bullet.Type = 1  # Bullet
                        para.ParagraphFormat.Bullet.Character = properties['bullet_character']
    
    # Apply animations if specified
    apply_animation_win32(shape, properties)
    
    return shape

def add_shape_element_win32(slide, content, properties):
    """Add a shape element to the slide using win32com.
    
    Args:
        slide: The PowerPoint slide object
        content: Text content for the shape
        properties: Dictionary of shape properties
    """
    # Get position and size
    left, top, width, height = get_position_and_size(properties)
    
    # Map shape type to PowerPoint MsoAutoShapeType
    shape_type = properties.get('shape_type', 'rectangle').lower()
    shape_map = {
        'rectangle': 1,
        'rounded_rectangle': 5,
        'oval': 9,
        'circle': 9,  # Same as oval
        'triangle': 6,
        'right_triangle': 8,
        'diamond': 4,
        'pentagon': 56,
        'hexagon': 10,
        'star': 12,
        'arrow': 13,
        'line': 20,
        'arc': 25,
        'cloud': 52,
        'heart': 74,
        'lightningbolt': 73,
        'sun': 23
    }
    shape_id = shape_map.get(shape_type, 1)  # Default to rectangle
    
    # Create shape
    shape = slide.Shapes.AddShape(shape_id, left, top, width, height)
    
    # Add text if present
    if content:
        shape.TextFrame.TextRange.Text = content
        
        # Word wrap for text in shapes
        shape.TextFrame.WordWrap = True
    
    # Set fill color
    if 'fill' in properties:
        fill_color = properties['fill']
        if fill_color.startswith('#'):
            shape.Fill.Visible = True
            shape.Fill.Solid()
            shape.Fill.ForeColor.RGB = hex_to_rgb(fill_color)
        elif fill_color.lower() == 'none' or fill_color.lower() == 'transparent':
            shape.Fill.Visible = False
    
    # Set border/line color
    if 'border_color' in properties:
        border_color = properties['border_color']
        if border_color.startswith('#'):
            shape.Line.Visible = True
            shape.Line.ForeColor.RGB = hex_to_rgb(border_color)
        elif border_color.lower() == 'none' or border_color.lower() == 'transparent':
            shape.Line.Visible = False
    
    # Set border/line width
    if 'border_width' in properties:
        width_match = re.match(r'(\d+(\.\d+)?)(pt|px|in|cm)?', properties['border_width'])
        if width_match:
            width_val = float(width_match.group(1))
            width_unit = width_match.group(3) or 'pt'
            
            # Convert to points
            if width_unit == 'in':
                width_pt = width_val * 72
            elif width_unit == 'cm':
                width_pt = width_val * 28.35
            elif width_unit == 'px':
                width_pt = width_val * 0.75
            else:  # pt
                width_pt = width_val
                
            shape.Line.Weight = width_pt
    
    # Set transparency
    if 'transparency' in properties:
        try:
            trans_val = float(properties['transparency'])
            # PowerPoint expects transparency as a percentage (0-1)
            if 0 <= trans_val <= 1:
                shape.Fill.Transparency = trans_val
        except ValueError:
            print(f"Warning: Invalid transparency value '{properties['transparency']}'")
    
    # Set rotation
    if 'rotation' in properties:
        try:
            rotation = float(properties['rotation'])
            shape.Rotation = rotation
        except ValueError:
            print(f"Warning: Invalid rotation value '{properties['rotation']}'")
    
    # Add shadow
    if 'shadow' in properties and properties['shadow'].lower() in ('true', 'yes', '1'):
        shape.Shadow.Visible = True
        
        # Shadow properties
        if 'shadow_color' in properties:
            shadow_color = properties['shadow_color']
            if shadow_color.startswith('#'):
                shape.Shadow.ForeColor.RGB = hex_to_rgb(shadow_color)
                
        if 'shadow_direction' in properties:
            try:
                direction = int(properties['shadow_direction'])
                shape.Shadow.OffsetX = 4 * round(math.cos(math.radians(direction)))
                shape.Shadow.OffsetY = 4 * round(math.sin(math.radians(direction)))
            except ValueError:
                print(f"Warning: Invalid shadow direction '{properties['shadow_direction']}'")
    
    # Text formatting (for shapes with text)
    if content:
        if 'font' in properties:
            shape.TextFrame.TextRange.Font.Name = properties['font']
        
        if 'font_size' in properties:
            try:
                size = float(properties['font_size'])
                shape.TextFrame.TextRange.Font.Size = size
            except ValueError:
                print(f"Warning: Invalid font size '{properties['font_size']}'")
        
        if 'font_color' in properties:
            color = properties['font_color']
            if color.startswith('#'):
                shape.TextFrame.TextRange.Font.Color.RGB = hex_to_rgb(color)
        
        if 'bold' in properties and properties['bold'].lower() in ('true', 'yes', '1'):
            shape.TextFrame.TextRange.Font.Bold = True
        
        if 'italic' in properties and properties['italic'].lower() in ('true', 'yes', '1'):
            shape.TextFrame.TextRange.Font.Italic = True
        
        if 'align' in properties:
            align = properties['align'].lower()
            if align == 'center':
                shape.TextFrame.TextRange.ParagraphFormat.Alignment = 2  # Center
            elif align == 'right':
                shape.TextFrame.TextRange.ParagraphFormat.Alignment = 3  # Right
            elif align == 'justify':
                shape.TextFrame.TextRange.ParagraphFormat.Alignment = 4  # Justify
        
        if 'vertical_align' in properties:
            valign = properties['vertical_align'].lower()
            if valign == 'middle':
                shape.TextFrame.VerticalAnchor = 3  # Middle
            elif valign == 'bottom':
                shape.TextFrame.VerticalAnchor = 2  # Bottom
    
    # Add hyperlink if specified
    if 'hyperlink' in properties:
        href = properties['hyperlink']
        tooltip = properties.get('tooltip', '')
        shape.ActionSettings(1).Hyperlink.Address = href  # 1 = mouse click
        shape.ActionSettings(1).Hyperlink.ScreenTip = tooltip
    
    # Apply animations if specified
    apply_animation_win32(shape, properties)
    
    return shape

def add_image_element_win32(slide, content, properties):
    """Add an image element to the slide using win32com.
    
    Args:
        slide: The PowerPoint slide object
        content: Path to the image file
        properties: Dictionary of image properties
    """
    # Get position and size
    left, top, width, height = get_position_and_size(properties)
    
    # Clean up image path
    image_path = content.strip()
    if not os.path.isabs(image_path):
        # Relative to the current directory
        image_path = os.path.abspath(image_path)
    
    # Check if file exists
    if not os.path.exists(image_path):
        print(f"Warning: Image file not found: {image_path}")
        return None
    
    # Add picture to slide
    shape = slide.Shapes.AddPicture(FileName=image_path, LinkToFile=False, 
                                  SaveWithDocument=True, Left=left, Top=top, 
                                  Width=width, Height=height)
    
    # Apply properties
    # Set rotation
    if 'rotation' in properties:
        try:
            rotation = float(properties['rotation'])
            shape.Rotation = rotation
        except ValueError:
            print(f"Warning: Invalid rotation value '{properties['rotation']}'")
    
    # Set transparency
    if 'transparency' in properties:
        try:
            trans_val = float(properties['transparency'])
            # PowerPoint PictureFormat doesn't support transparency directly
            # For Win32com, we can use ColorFormat.TintAndShade to approximate
            shape.PictureFormat.Brightness = 1.0 - trans_val
        except ValueError:
            print(f"Warning: Invalid transparency value '{properties['transparency']}'")
    
    # Set border
    if 'border_color' in properties:
        border_color = properties['border_color']
        if border_color.startswith('#'):
            shape.Line.Visible = True
            shape.Line.ForeColor.RGB = hex_to_rgb(border_color)
    
    if 'border_width' in properties:
        width_match = re.match(r'(\d+(\.\d+)?)(pt|px|in|cm)?', properties['border_width'])
        if width_match:
            width_val = float(width_match.group(1))
            width_unit = width_match.group(3) or 'pt'
            
            # Convert to points
            if width_unit == 'in':
                width_pt = width_val * 72
            elif width_unit == 'cm':
                width_pt = width_val * 28.35
            elif width_unit == 'px':
                width_pt = width_val * 0.75
            else:  # pt
                width_pt = width_val
                
            shape.Line.Weight = width_pt
    
    # Apply image adjustments
    # Brightness
    if 'brightness' in properties:
        try:
            brightness = float(properties['brightness'])
            # PowerPoint expects brightness from -1 to 1
            if -1 <= brightness <= 1:
                shape.PictureFormat.Brightness = brightness
        except ValueError:
            print(f"Warning: Invalid brightness value '{properties['brightness']}'")
    
    # Contrast
    if 'contrast' in properties:
        try:
            contrast = float(properties['contrast'])
            # PowerPoint expects contrast from -1 to 1
            if -1 <= contrast <= 1:
                shape.PictureFormat.Contrast = contrast
        except ValueError:
            print(f"Warning: Invalid contrast value '{properties['contrast']}'")
    
    # Apply animations if specified
    apply_animation_win32(shape, properties)
    
    return shape

def add_table_element_win32(slide, content, properties):
    """Add a table element to the slide using win32com.
    
    Args:
        slide: The PowerPoint slide object
        content: The table content as string
        properties: Dictionary of table properties
    """
    # Get position and size
    left, top, width, height = get_position_and_size(properties)
    
    # Parse table content
    rows = []
    
    # Handle markdown table format with pipes
    if '|' in content:
        for line in content.strip().split('\n'):
            if line.strip() and not line.startswith('|-'):  # Skip separator rows
                cells = [cell.strip() for cell in line.split('|')]
                # Remove empty cells at beginning and end caused by leading/trailing |
                if not cells[0]:
                    cells = cells[1:]
                if not cells[-1]:
                    cells = cells[:-1]
                rows.append(cells)
    # Handle CSV-style format
    else:
        for line in content.strip().split('\n'):
            if line.strip():
                cells = [cell.strip() for cell in line.split(',')]
                rows.append(cells)
    
    if not rows:
        print("Warning: No data found for table")
        return
    
    # Determine table dimensions
    num_rows = len(rows)
    num_cols = max(len(row) for row in rows)
    
    # Create table
    table_shape = slide.Shapes.AddTable(num_rows, num_cols, left, top, width, height)
    table = table_shape.Table
    
    # Fill in the data
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < num_cols:  # Avoid index errors
                table.Cell(row_idx + 1, col_idx + 1).Shape.TextFrame.TextRange.Text = cell_text
    
    # Apply table style properties
    # Header row
    if properties.get('header_row', '').lower() in ('true', 'yes', '1'):
        table.FirstRow = True
        
        # Apply header formatting
        for col_idx in range(1, num_cols + 1):
            cell = table.Cell(1, col_idx)
            # Bold text
            cell.Shape.TextFrame.TextRange.Font.Bold = True
            # Background color for header (light gray)
            cell.Shape.Fill.ForeColor.RGB = 0xE0E0E0  # Light gray
    
    # First column formatting
    if properties.get('first_column', '').lower() in ('true', 'yes', '1'):
        table.FirstColumn = True
        for row_idx in range(1, num_rows + 1):
            cell = table.Cell(row_idx, 1)
            cell.Shape.TextFrame.TextRange.Font.Bold = True
    
    # Banded rows
    if properties.get('banded_rows', '').lower() in ('true', 'yes', '1'):
        table.BandRows = True
    
    # Banded columns
    if properties.get('banded_columns', '').lower() in ('true', 'yes', '1'):
        table.BandColumns = True
    
    # Border color
    if 'border_color' in properties:
        color = properties['border_color']
        if color.startswith('#'):
            rgb = hex_to_rgb(color)
            # Apply to all borders
            table.Borders.ForeColor.RGB = rgb
    
    # Border width
    if 'border_width' in properties:
        # Convert width to points
        width_match = re.match(r'(\d+(\.\d+)?)(pt|px|in|cm)?', properties['border_width'])
        if width_match:
            width_val = float(width_match.group(1))
            width_unit = width_match.group(3) or 'pt'
            
            # Convert to points
            if width_unit == 'in':
                width_pt = width_val * 72
            elif width_unit == 'cm':
                width_pt = width_val * 28.35
            elif width_unit == 'px':
                width_pt = width_val * 0.75
            else:  # pt
                width_pt = width_val
                
            # Apply width to all borders
            for border in [table.Borders.InsideHorizontal, 
                          table.Borders.InsideVertical, 
                          table.Borders.OutsideTop,
                          table.Borders.OutsideBottom,
                          table.Borders.OutsideLeft,
                          table.Borders.OutsideRight]:
                border.Weight = width_pt
    
    # Table style
    if 'table_style' in properties:
        style_name = properties['table_style'].lower()
        style_map = {
            'light': 15,   # Light Style 1
            'medium': 4,   # Medium Style 1
            'dark': 13     # Dark Style 1
        }
        
        if style_name in style_map:
            table.ApplyStyle(style_map[style_name])
    
    # Cell padding
    if 'cell_padding' in properties:
        padding_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', properties['cell_padding'])
        if padding_match:
            padding_val = float(padding_match.group(1))
            padding_unit = padding_match.group(3) or 'in'
            
            # Convert to points
            if padding_unit == 'in':
                padding_pt = padding_val * 72
            elif padding_unit == 'cm':
                padding_pt = padding_val * 28.35
            elif padding_unit == 'px':
                padding_pt = padding_val * 0.75
            else:  # pt
                padding_pt = padding_val
                
            # Apply padding
            for row_idx in range(1, num_rows + 1):
                for col_idx in range(1, num_cols + 1):
                    cell = table.Cell(row_idx, col_idx)
                    cell.Shape.TextFrame.MarginLeft = padding_pt
                    cell.Shape.TextFrame.MarginRight = padding_pt
                    cell.Shape.TextFrame.MarginTop = padding_pt
                    cell.Shape.TextFrame.MarginBottom = padding_pt
    
    # Apply animations if specified
    apply_animation_win32(table_shape, properties)
    
    return table_shape

def add_chart_element_win32(slide, content, properties):
    """Add a chart element to the slide using win32com.
    
    Args:
        slide: The PowerPoint slide object
        content: The chart data as string
        properties: Dictionary of chart properties
    """
    import win32com.client
    
    # Get position and size
    left, top, width, height = get_position_and_size(properties)
    
    # Parse chart data
    data_rows = []
    for line in content.strip().split('\n'):
        if line.strip():
            cells = [cell.strip() for cell in line.split(',')]
            data_rows.append(cells)
    
    if not data_rows:
        print("Warning: No data found for chart")
        return
    
    # Determine chart type
    chart_type = properties.get('chart_type', 'column').lower()
    chart_type_map = {
        'column': win32com.client.constants.xlColumnClustered,
        'bar': win32com.client.constants.xlBarClustered,
        'line': win32com.client.constants.xlLine,
        'pie': win32com.client.constants.xlPie,
        'area': win32com.client.constants.xlAreaStacked,
        'scatter': win32com.client.constants.xlXYScatter,
        'doughnut': win32com.client.constants.xlDoughnut,
        'radar': win32com.client.constants.xlRadar
    }
    chart_type_id = chart_type_map.get(chart_type, win32com.client.constants.xlColumnClustered)
    
    # Create chart
    chart_shape = slide.Shapes.AddChart2(-1, chart_type_id, left, top, width, height)
    chart = chart_shape.Chart
    
    # Create Excel workbook for chart data
    workbook = chart.ChartData.Workbook
    worksheet = workbook.Worksheets(1)
    
    # Fill in the data
    num_rows = len(data_rows)
    num_cols = max(len(row) for row in data_rows)
    
    for row_idx, row_data in enumerate(data_rows):
        for col_idx, cell_value in enumerate(row_data):
            cell = worksheet.Cells(row_idx + 1, col_idx + 1)
            
            # Try to convert to number if possible
            try:
                if '.' in cell_value:
                    cell.Value = float(cell_value)
                else:
                    cell.Value = int(cell_value)
            except ValueError:
                cell.Value = cell_value
    
    # Set chart range
    chart.SetSourceData(worksheet.Range(
        worksheet.Cells(1, 1),
        worksheet.Cells(num_rows, num_cols)
    ))
    
    # Apply chart properties
    # Chart title
    if 'title' in properties:
        chart.HasTitle = True
        chart.ChartTitle.Text = properties['title']
    
    # Axis titles
    if 'x_axis_title' in properties and chart_type != 'pie' and chart_type != 'doughnut':
        chart.Axes(1).HasTitle = True  # 1 = x-axis
        chart.Axes(1).AxisTitle.Text = properties['x_axis_title']
    
    if 'y_axis_title' in properties and chart_type != 'pie' and chart_type != 'doughnut':
        chart.Axes(2).HasTitle = True  # 2 = y-axis
        chart.Axes(2).AxisTitle.Text = properties['y_axis_title']
    
    # Legend
    has_legend = properties.get('has_legend', '').lower() in ('true', 'yes', '1')
    chart.HasLegend = has_legend
    
    if has_legend and 'legend_position' in properties:
        position = properties['legend_position'].lower()
        position_map = {
            'right': 2,
            'left': 3,
            'top': 1,
            'bottom': 4
        }
        if position in position_map:
            chart.Legend.Position = position_map[position]
    
    # Data labels
    if properties.get('data_labels', '').lower() in ('true', 'yes', '1'):
        # This is a simplification - different chart types have different ways to show data labels
        try:
            chart.ApplyDataLabels()
        except:
            print("Warning: Could not apply data labels to chart")
    
    # Gridlines
    if properties.get('gridlines', '').lower() in ('true', 'no', '0'):
        # Disable gridlines (they're on by default)
        try:
            chart.Axes(1).HasMajorGridlines = False
            chart.Axes(2).HasMajorGridlines = False
        except:
            print("Warning: Could not disable gridlines")
    
    # Apply animations if specified
    apply_animation_win32(chart_shape, properties)
    
    return chart_shape

def add_smartart_element_win32(slide, content, properties):
    """Add a SmartArt element to the slide using win32com.
    
    Args:
        slide: The PowerPoint slide object
        content: The SmartArt data as string
        properties: Dictionary of SmartArt properties
    """
    # Get position and size
    left, top, width, height = get_position_and_size(properties)
    
    # Parse SmartArt data (comma-separated list of text items)
    items = [item.strip() for item in content.strip().split(',')]
    
    if not items:
        print("Warning: No data found for SmartArt")
        return None
    
    try:
        # Determine SmartArt type
        smartart_type = properties.get('type', 'process').lower()
        
        # Map from our simple type names to MS Office SmartArt layout names
        layout_map = {
            'process': 'BasicProcess',
            'cycle': 'BasicCycle',
            'hierarchy': 'Hierarchy',
            'pyramid': 'BasicPyramid',
            'radial': 'BasicRadial',
            'venn': 'BasicVenn',
            'matrix': 'BasicMatrix',
            'relationship': 'BasicChevron',
            'list': 'BasicBlockList'
        }
        
        layout_name = layout_map.get(smartart_type, 'BasicProcess')
        print(f"  Creating SmartArt with layout: {layout_name}")
        
        # Create SmartArt
        smart_art = slide.Shapes.AddSmartArt(layout_name, left, top, width, height)
        
        # Add text to SmartArt nodes
        nodes = smart_art.SmartArt.AllNodes
        
        print(f"  SmartArt has {nodes.Count} nodes, adding {len(items)} items")
        
        # Fill in as many nodes as we have data for
        for i, item in enumerate(items):
            if i < nodes.Count:
                print(f"  Setting node {i+1} text to: {item}")
                nodes.Item(i + 1).TextFrame2.TextRange.Text = item
        
        # Apply animations if specified
        apply_animation_win32(smart_art, properties)
        
        return smart_art
    except Exception as e:
        print(f"Error creating SmartArt: {e}")
        print("Falling back to text boxes for process visualization")
        
        # Create a fallback visualization using shapes and text
        box_width = width / len(items)
        for i, item in enumerate(items):
            box_left = left + (i * box_width)
            # Create a box
            box = slide.Shapes.AddShape(1, box_left, top, box_width * 0.9, height * 0.8)
            box.Fill.Solid()
            box.Fill.ForeColor.RGB = hex_to_rgb("#4472C4")  # Blue
            box.Line.ForeColor.RGB = hex_to_rgb("#2F528F")  # Darker blue
            
            # Add text
            box.TextFrame.TextRange.Text = item
            box.TextFrame.TextRange.Font.Color.RGB = hex_to_rgb("#FFFFFF")  # White
            box.TextFrame.TextRange.ParagraphFormat.Alignment = 2  # Center
            box.TextFrame.VerticalAnchor = 3  # Middle
            
            # Add connector to next box if not the last one
            if i < len(items) - 1:
                connector = slide.Shapes.AddLine(
                    box_left + box_width * 0.9, top + height * 0.4,
                    box_left + box_width * 1.1, top + height * 0.4
                )
                connector.Line.ForeColor.RGB = hex_to_rgb("#2F528F")
                connector.Line.Weight = 2
                
                # Add arrowhead
                connector.Line.EndArrowheadStyle = 5  # Triangle arrowhead
        
        return None

def apply_animation_win32(shape, properties):
    """Apply animation effects to a shape using win32com."""
    # For testing, comment this out to disable animations
    if 'animation' not in properties:
        return

    try:
        # Get the slide containing the shape
        slide = shape.Parent
        
        # Get the animation sequence
        sequence = slide.TimeLine.MainSequence
        
        # Map animation types to PowerPoint constants
        animation_map = {
            'fade': 4,                # ppAnimEffectFade
            'appear': 3,              # ppAnimEffectAppear
            'fly_in': 11,             # ppAnimEffectFly
            'float': 16,              # ppAnimEffectFloat
            'split': 13,              # ppAnimEffectSplit
            'wipe': 14,               # ppAnimEffectWipe
            'zoom': 25,               # ppAnimEffectGrowShrink
            'bounce': 27,             # ppAnimEffectBounce
            'spin': 41,               # ppAnimEffectSpin
            'swivel': 40,             # ppAnimEffectSwivel
            'pulse': 38,              # ppAnimEffectPulse
            'color': 42,              # ppAnimEffectTeeter
            'grow': 25                # ppAnimEffectGrowShrink
        }
        
        # Get animation effect
        animation_type = properties['animation'].lower()
        effect_id = animation_map.get(animation_type, 4)  # Default to fade
        
        # Add animation effect
        effect = sequence.AddEffect(shape, effect_id, 0, 1)  # 0 = entrance, 1 = after previous
        
        # Animation timing
        if 'animation_trigger' in properties:
            trigger = properties['animation_trigger'].lower()
            if trigger == 'on_click':
                effect.Timing.TriggerType = 1  # On click
            elif trigger == 'with_previous':
                effect.Timing.TriggerType = 2  # With previous
            elif trigger == 'after_previous':
                effect.Timing.TriggerType = 3  # After previous
        
        print(f"  Added animation: {animation_type}")
        return effect
    except Exception as e:
        print(f"Error applying animation: {e}")
        return None

def create_presentation_win32(slides, settings, output_file):
    """Create a PowerPoint presentation using win32com.
    
    Args:
        slides: List of slide data dictionaries
        settings: Dictionary of global settings
        output_file: Path to save the output file
        
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        import win32com.client
        
        # Create PowerPoint application
        print("Starting PowerPoint...")
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = True  # Make it visible for debugging
        
        # Create new presentation or use template
        print("Creating presentation...")
        if 'template' in settings and os.path.exists(settings['template']):
            template_path = os.path.abspath(settings['template'])
            presentation = ppt.Presentations.Open(template_path, WithWindow=False)
            # Clear existing slides if needed
            while presentation.Slides.Count > 0:
                presentation.Slides(1).Delete()
        else:
            presentation = ppt.Presentations.Add()
        
        # Apply global settings
        if 'slide_width' in settings and 'slide_height' in settings:
            try:
                # Convert to points (1 inch = 72 points)
                w_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(settings['slide_width']))
                h_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(settings['slide_height']))
                
                if w_match and h_match:
                    w_num = float(w_match.group(1))
                    h_num = float(h_match.group(1))
                    
                    w_unit = w_match.group(3) or 'in'
                    h_unit = h_match.group(3) or 'in'
                    
                    # Convert to points
                    if w_unit == 'in': w_points = w_num * 72
                    elif w_unit == 'cm': w_points = w_num * 28.35
                    elif w_unit == 'pt': w_points = w_num
                    elif w_unit == 'px': w_points = w_num * 0.75
                    
                    if h_unit == 'in': h_points = h_num * 72
                    elif h_unit == 'cm': h_points = h_num * 28.35
                    elif h_unit == 'pt': h_points = h_num
                    elif h_unit == 'px': h_points = h_num * 0.75
                    
                    presentation.PageSetup.SlideWidth = w_points
                    presentation.PageSetup.SlideHeight = h_points
            except Exception as e:
                print(f"Warning: Could not set slide dimensions: {e}")
        
        # Check for company logo in global settings
        company_logo = None
        logo_path = settings.get('company_logo', '')
        if logo_path and os.path.exists(logo_path):
            company_logo = os.path.abspath(logo_path)
        
        # Process each slide
        print(f"Adding {len(slides)} slides...")
        for i, slide_data in enumerate(slides):
            # Get layout
            layout_name = slide_data['settings'].get('layout', 'Blank').lower()
            layout_idx = 6  # Default to blank
            
            if layout_name == 'title slide': layout_idx = 0
            elif layout_name == 'title and content': layout_idx = 1
            elif layout_name == 'section header': layout_idx = 2
            elif layout_name == 'two content': layout_idx = 3
            elif layout_name == 'comparison': layout_idx = 4
            elif layout_name == 'title only': layout_idx = 5
            elif layout_name == 'blank': layout_idx = 6
            elif layout_name == 'content with caption': layout_idx = 7
            elif layout_name == 'picture with caption': layout_idx = 8
            
            # Add slide
            slide = presentation.Slides.Add(i+1, layout_idx+1)
            print(f"  Slide {i+1}: {slide_data['title']} - Layout: {layout_name}")
            
            # Set title if layout has title
            if slide.Shapes.HasTitle and slide_data['title']:
                slide.Shapes.Title.TextFrame.TextRange.Text = slide_data['title']
            
            # Set background
            if 'background' in slide_data['settings']:
                bg_value = slide_data['settings']['background']
                
                # Check if it's a color (starts with #) or a file path
                if bg_value.startswith('#'):
                    # It's a color
                    rgb = hex_to_rgb(bg_value)
                    slide.Background.Fill.Visible = True
                    slide.Background.Fill.Solid()
                    slide.Background.Fill.ForeColor.RGB = rgb
                elif os.path.exists(bg_value):
                    # It's a file
                    img_path = os.path.abspath(bg_value)
                    slide.Background.Fill.Visible = True
                    slide.Background.Fill.UserPicture(img_path)
            
            # Set slide transition
            if 'transition' in slide_data['settings']:
                transition_name = slide_data['settings']['transition'].lower()
                transition_map = {
                    'none': 0,
                    'fade': 4,
                    'push': 6,
                    'wipe': 3,
                    'split': 2,
                    'reveal': 1,
                    'zoom': 24,
                    'dissolve': 5
                }
                
                transition_id = transition_map.get(transition_name, 0)
                if transition_id > 0:
                    slide.SlideShowTransition.EntryEffect = transition_id
                    
                    # Transition speed
                    if 'transition_speed' in slide_data['settings']:
                        speed = slide_data['settings']['transition_speed'].lower()
                        if speed == 'slow':
                            slide.SlideShowTransition.Speed = 3
                        elif speed == 'medium':
                            slide.SlideShowTransition.Speed = 2
                        elif speed == 'fast':
                            slide.SlideShowTransition.Speed = 1
            
            # Process elements in order
            
            # Process elements in order
            for element in slide_data['elements']:
                element_type = element['type'].lower()
                properties = element['properties']
                content = element['content']
                
                # Add debug info
                debug_element(element_type, properties, content)
                
                try:
                    if element_type == 'text':
                        add_text_element_win32(slide, content, properties)
                    
                    elif element_type == 'shape':
                        add_shape_element_win32(slide, content, properties)
                    
                    elif element_type == 'image':
                        add_image_element_win32(slide, content, properties)
                    
                    elif element_type == 'table':
                        add_table_element_win32(slide, content, properties)
                    
                    elif element_type == 'chart':
                        add_chart_element_win32(slide, content, properties)
                        
                    elif element_type == 'smartart':
                        add_smartart_element_win32(slide, content, properties)
                    
                    else:
                        print(f"Warning: Unknown element type '{element_type}'")
                except Exception as e:
                    print(f"Error adding {element_type} element: {e}")
                    traceback.print_exc()
            
            # Add company logo if specified in global settings and not already on the slide
            if company_logo:
                # Check if footer is enabled for this slide
                footer_enabled = slide_data['settings'].get('footer', '').lower() in ('true', 'yes', '1')
                
                # Add logo in the bottom right if footer is enabled
                if footer_enabled:
                    logo_width = 72  # 1 inch
                    logo_height = 36  # 0.5 inch
                    logo_left = presentation.PageSetup.SlideWidth - logo_width - 36  # 0.5 inch from right
                    logo_top = presentation.PageSetup.SlideHeight - logo_height - 36  # 0.5 inch from bottom
                    
                    slide.Shapes.AddPicture(
                        FileName=company_logo,
                        LinkToFile=False,
                        SaveWithDocument=True,
                        Left=logo_left,
                        Top=logo_top,
                        Width=logo_width,
                        Height=logo_height
                    )
                    
                    # Add footer text if specified
                    if 'footer_text' in settings:
                        footer_text = settings['footer_text']
                        footer_left = 36  # 0.5 inch from left
                        footer_top = presentation.PageSetup.SlideHeight - 36  # 0.5 inch from bottom
                        footer_width = presentation.PageSetup.SlideWidth - 144  # Subtract logo width and margins
                        
                        footer_shape = slide.Shapes.AddTextbox(
                            1, footer_left, footer_top - 18, footer_width, 36
                        )
                        footer_shape.TextFrame.TextRange.Text = footer_text
                        footer_shape.TextFrame.TextRange.Font.Size = 10
                        footer_shape.TextFrame.TextRange.Font.Color.RGB = 0x808080  # Gray
            
            # Speaker notes
            if 'notes' in slide_data['settings']:
                notes_text = slide_data['settings']['notes']
                slide.NotesPage.Shapes.Placeholders(2).TextFrame.TextRange.Text = notes_text
        
        # Save the presentation
        output_path = os.path.abspath(output_file)
        print(f"Saving presentation to: {output_path}")
        presentation.SaveAs(output_path)
        print(f"Presentation saved successfully to: {output_path}")
        
        return True
        
    except Exception as e:
        print(f"\n✗ ERROR: {e}")
        print("\nDetailed error information:")
        traceback.print_exc()
        print("\nPossible solutions:")
        print("1. Make sure Microsoft PowerPoint is installed")
        print("2. Try running: pip install --upgrade pywin32")
        print("3. If using a virtual environment, make sure to run: python -m pywin32_postinstall -install")
        return False

def create_presentation_python_pptx(slides, settings, output_file):
    """Create a PowerPoint presentation using python-pptx library.
    
    Args:
        slides: List of slide data dictionaries
        settings: Dictionary of global settings
        output_file: Path to save the output file
        
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Cm, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
        from pptx.dml.color import RGBColor
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        
        # Create presentation
        print("Creating presentation using python-pptx...")
        prs = Presentation()
        
        # Apply global settings
        if 'slide_width' in settings and 'slide_height' in settings:
            try:
                # Convert to inches (default unit for python-pptx)
                w_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(settings['slide_width']))
                h_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(settings['slide_height']))
                
                if w_match and h_match:
                    w_num = float(w_match.group(1))
                    h_num = float(h_match.group(1))
                    
                    w_unit = w_match.group(3) or 'in'
                    h_unit = h_match.group(3) or 'in'
                    
                    # Convert to appropriate pptx units
                    if w_unit == 'in': w_val = Inches(w_num)
                    elif w_unit == 'cm': w_val = Cm(w_num)
                    elif w_unit == 'pt': w_val = Pt(w_num)
                    elif w_unit == 'px': w_val = Inches(w_num / 96)  # Approximate conversion
                    
                    if h_unit == 'in': h_val = Inches(h_num)
                    elif h_unit == 'cm': h_val = Cm(h_num)
                    elif h_unit == 'pt': h_val = Pt(h_num)
                    elif h_unit == 'px': h_val = Inches(h_num / 96)  # Approximate conversion
                    
                    prs.slide_width = w_val
                    prs.slide_height = h_val
            except Exception as e:
                print(f"Warning: Could not set slide dimensions: {e}")
        
        # Process each slide
        print(f"Adding {len(slides)} slides...")
        for i, slide_data in enumerate(slides):
            # Get layout
            layout_name = slide_data['settings'].get('layout', 'Blank').lower()
            
            # Map layout name to slide layout in presentation
            layout_idx = 6  # Default to blank
            if layout_name == 'title slide': layout_idx = 0
            elif layout_name == 'title and content': layout_idx = 1
            elif layout_name == 'section header': layout_idx = 2
            elif layout_name == 'two content': layout_idx = 3
            elif layout_name == 'comparison': layout_idx = 4
            elif layout_name == 'title only': layout_idx = 5
            elif layout_name == 'blank': layout_idx = 6
            
            # Ensure layout_idx is in range
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = min(len(prs.slide_layouts) - 1, layout_idx)
            
            slide_layout = prs.slide_layouts[layout_idx]
            
            # Add slide
            slide = prs.slides.add_slide(slide_layout)
            print(f"  Slide {i+1}: {slide_data['title']} - Layout: {layout_name}")
            
            # Set title if layout has title
            if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
                if slide_data['title']:
                    slide.shapes.title.text = slide_data['title']
            
            # Note: Background settings require different handling in python-pptx
            # Set slide background if specified
            if 'background' in slide_data['settings']:
                bg_value = slide_data['settings']['background']
                
                # For python-pptx, setting slide backgrounds is more complex
                # and requires using the slide.background.fill properties
                # Implementation would go here, but is more complex than space allows
                print(f"  Note: Setting custom background color '{bg_value}' requires win32com backend")
            
            # Process elements
            for element in slide_data['elements']:
                element_type = element['type'].lower()
                properties = element['properties']
                content = element['content']
                
                # Parse position and size (converting to pptx units)
                left, top, width, height = 0, 0, 0, 0
                
                if 'x' in properties:
                    x_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(properties['x']))
                    if x_match:
                        x_num = float(x_match.group(1))
                        x_unit = x_match.group(3) or 'in'
                        if x_unit == 'in': left = Inches(x_num)
                        elif x_unit == 'cm': left = Cm(x_num)
                        elif x_unit == 'pt': left = Pt(x_num)
                        elif x_unit == 'px': left = Inches(x_num / 96)
                
                if 'y' in properties:
                    y_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(properties['y']))
                    if y_match:
                        y_num = float(y_match.group(1))
                        y_unit = y_match.group(3) or 'in'
                        if y_unit == 'in': top = Inches(y_num)
                        elif y_unit == 'cm': top = Cm(y_num)
                        elif y_unit == 'pt': top = Pt(y_num)
                        elif y_unit == 'px': top = Inches(y_num / 96)
                
                if 'width' in properties:
                    w_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(properties['width']))
                    if w_match:
                        w_num = float(w_match.group(1))
                        w_unit = w_match.group(3) or 'in'
                        if w_unit == 'in': width = Inches(w_num)
                        elif w_unit == 'cm': width = Cm(w_num)
                        elif w_unit == 'pt': width = Pt(w_num)
                        elif w_unit == 'px': width = Inches(w_num / 96)
                
                if 'height' in properties:
                    h_match = re.match(r'(\d+(\.\d+)?)(in|cm|pt|px)?', str(properties['height']))
                    if h_match:
                        h_num = float(h_match.group(1))
                        h_unit = h_match.group(3) or 'in'
                        if h_unit == 'in': height = Inches(h_num)
                        elif h_unit == 'cm': height = Cm(h_num)
                        elif h_unit == 'pt': height = Pt(h_num)
                        elif h_unit == 'px': height = Inches(h_num / 96)
                
                # Add appropriate element based on type
                if element_type == 'text':
                    # Add text box
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = content
                    
                    # Apply text formatting
                    p = text_frame.paragraphs[0]
                    
                    if 'font_size' in properties:
                        try:
                            size = float(properties['font_size'])
                            p.font.size = Pt(size)
                        except ValueError:
                            print(f"Warning: Invalid font size '{properties['font_size']}'")
                    
                    if 'font_color' in properties:
                        color = properties['font_color']
                        if color.startswith('#'):
                            hex_color = color.lstrip('#')
                            if len(hex_color) == 6:
                                r = int(hex_color[0:2], 16)
                                g = int(hex_color[2:4], 16)
                                b = int(hex_color[4:6], 16)
                                p.font.color.rgb = RGBColor(r, g, b)
                    
                    if 'bold' in properties and properties['bold'].lower() in ('true', 'yes', '1'):
                        p.font.bold = True
                    
                    if 'italic' in properties and properties['italic'].lower() in ('true', 'yes', '1'):
                        p.font.italic = True
                    
                    if 'align' in properties:
                        align = properties['align'].lower()
                        if align == 'center':
                            p.alignment = PP_ALIGN.CENTER
                        elif align == 'right':
                            p.alignment = PP_ALIGN.RIGHT
                        elif align == 'justify':
                            p.alignment = PP_ALIGN.JUSTIFY
                
                elif element_type == 'shape':
                    # Map shape type to pptx AutoShape enum
                    shape_type = properties.get('shape_type', 'rectangle').lower()
                    shape_map = {
                        'rectangle': MSO_SHAPE.RECTANGLE,
                        'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
                        'oval': MSO_SHAPE.OVAL,
                        'diamond': MSO_SHAPE.DIAMOND,
                        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
                        'right_triangle': MSO_SHAPE.RIGHT_TRIANGLE,
                        'pentagon': MSO_SHAPE.PENTAGON,
                        'hexagon': MSO_SHAPE.HEXAGON,
                        'star': MSO_SHAPE.STAR_5_POINTS,
                        'arrow': MSO_SHAPE.RIGHT_ARROW
                    }
                    shape_id = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)
                    
                    # Create shape
                    shape = slide.shapes.add_shape(shape_id, left, top, width, height)
                    
                    # Set fill color
                    if 'fill' in properties:
                        fill_color = properties['fill']
                        if fill_color.startswith('#'):
                            hex_color = fill_color.lstrip('#')
                            if len(hex_color) == 6:
                                r = int(hex_color[0:2], 16)
                                g = int(hex_color[2:4], 16)
                                b = int(hex_color[4:6], 16)
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = RGBColor(r, g, b)
                        elif fill_color.lower() == 'none' or fill_color.lower() == 'transparent':
                            shape.fill.background()
                    
                    # Set line/border color
                    if 'border_color' in properties:
                        border_color = properties['border_color']
                        if border_color.startswith('#'):
                            hex_color = border_color.lstrip('#')
                            if len(hex_color) == 6:
                                r = int(hex_color[0:2], 16)
                                g = int(hex_color[2:4], 16)
                                b = int(hex_color[4:6], 16)
                                shape.line.color.rgb = RGBColor(r, g, b)
                    
                    # Set text
                    if content:
                        text_frame = shape.text_frame
                        text_frame.text = content
                        
                        # Text formatting
                        p = text_frame.paragraphs[0]
                        
                        if 'font_size' in properties:
                            try:
                                size = float(properties['font_size'])
                                p.font.size = Pt(size)
                            except ValueError:
                                print(f"Warning: Invalid font size '{properties['font_size']}'")
                        
                        if 'font_color' in properties:
                            color = properties['font_color']
                            if color.startswith('#'):
                                hex_color = color.lstrip('#')
                                if len(hex_color) == 6:
                                    r = int(hex_color[0:2], 16)
                                    g = int(hex_color[2:4], 16)
                                    b = int(hex_color[4:6], 16)
                                    p.font.color.rgb = RGBColor(r, g, b)
                        
                        if 'align' in properties:
                            align = properties['align'].lower()
                            if align == 'center':
                                p.alignment = PP_ALIGN.CENTER
                            elif align == 'right':
                                p.alignment = PP_ALIGN.RIGHT
                            elif align == 'justify':
                                p.alignment = PP_ALIGN.JUSTIFY
                        
                        if 'vertical_align' in properties:
                            valign = properties['vertical_align'].lower()
                            if valign == 'middle':
                                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                            elif valign == 'bottom':
                                text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
                
                elif element_type == 'image':
                    # Clean up image path
                    image_path = content.strip()
                    if not os.path.isabs(image_path):
                        # Relative to the current directory
                        image_path = os.path.abspath(image_path)
                    
                    # Check if file exists
                    if not os.path.exists(image_path):
                        print(f"Warning: Image file not found: {image_path}")
                        continue
                    
                    # Add picture
                    slide.shapes.add_picture(image_path, left, top, width, height)
                
                elif element_type == 'table':
                    # Parse table content
                    rows = []
                    
                    # Handle markdown table format with pipes
                    if '|' in content:
                        for line in content.strip().split('\n'):
                            if line.strip() and not line.startswith('|-'):  # Skip separator rows
                                cells = [cell.strip() for cell in line.split('|')]
                                # Remove empty cells at beginning and end caused by leading/trailing |
                                if not cells[0]:
                                    cells = cells[1:]
                                if not cells[-1]:
                                    cells = cells[:-1]
                                rows.append(cells)
                    # Handle CSV-style format
                    else:
                        for line in content.strip().split('\n'):
                            if line.strip():
                                cells = [cell.strip() for cell in line.split(',')]
                                rows.append(cells)
                    
                    if not rows:
                        print("Warning: No data found for table")
                        continue
                    
                    # Determine table dimensions
                    num_rows = len(rows)
                    num_cols = max(len(row) for row in rows)
                    
                    # Create table
                    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
                    
                    # Fill in the data
                    for row_idx, row_data in enumerate(rows):
                        for col_idx, cell_text in enumerate(row_data):
                            if col_idx < num_cols:  # Avoid index errors
                                cell = table.cell(row_idx, col_idx)
                                cell.text = cell_text
                                
                                # Apply header formatting to first row if enabled
                                if row_idx == 0 and properties.get('header_row', '').lower() in ('true', 'yes', '1'):
                                    cell.text_frame.paragraphs[0].font.bold = True
                
                elif element_type == 'chart':
                    # Parse chart data
                    data_rows = []
                    for line in content.strip().split('\n'):
                        if line.strip():
                            cells = [cell.strip() for cell in line.split(',')]
                            data_rows.append(cells)
                    
                    if not data_rows:
                        print("Warning: No data found for chart")
                        continue
                    
                    # Determine chart type
                    chart_type = properties.get('chart_type', 'column').lower()
                    chart_type_map = {
                        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                        'line': XL_CHART_TYPE.LINE,
                        'pie': XL_CHART_TYPE.PIE,
                        'area': XL_CHART_TYPE.AREA,
                        'scatter': XL_CHART_TYPE.XY_SCATTER,
                        'doughnut': XL_CHART_TYPE.DOUGHNUT
                    }
                    chart_type_id = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
                    
                    # Prepare chart data
                    chart_data = CategoryChartData()
                    
                    # Add categories (first column)
                    categories = [row[0] for row in data_rows[1:]]  # Skip header row
                    chart_data.categories = categories
                    
                    # Add series (remaining columns)
                    for col_idx in range(1, len(data_rows[0])):
                        series_name = data_rows[0][col_idx]
                        series_values = []
                        
                        for row_idx in range(1, len(data_rows)):
                            if col_idx < len(data_rows[row_idx]):
                                try:
                                    val = float(data_rows[row_idx][col_idx])
                                    series_values.append(val)
                                except ValueError:
                                    series_values.append(0)
                                    print(f"Warning: Non-numeric value in chart data: {data_rows[row_idx][col_idx]}")
                        
                        chart_data.add_series(series_name, series_values)
                    
                    # Add chart to slide
                    chart = slide.shapes.add_chart(chart_type_id, left, top, width, height, chart_data).chart
                    
                    # Set chart title
                    if 'title' in properties:
                        chart.has_title = True
                        chart.chart_title.text_frame.text = properties['title']
        
        # Save the presentation
        output_path = os.path.abspath(output_file)
        print(f"Saving presentation to: {output_path}")
        prs.save(output_path)
        print(f"Presentation saved successfully to: {output_path}")
        
        return True
        
    except Exception as e:
        print(f"\n✗ ERROR: {e}")
        print("\nDetailed error information:")
        traceback.print_exc()
        print("\nPossible solutions:")
        print("1. Make sure python-pptx is installed: pip install python-pptx")
        print("2. Use the win32com backend for more advanced features")
        return False

def create_presentation(slides, settings, output_file, backend=BackendType.WIN32COM):
    """Create PowerPoint presentation from slide data using the specified backend.
    
    Args:
        slides: List of slide data dictionaries
        settings: Dictionary of global settings
        output_file: Path to save the output file
        backend: Backend type to use (win32com or python-pptx)
        
    Returns:
        bool: True if successful, False otherwise
    """
    if backend == BackendType.WIN32COM:
        return create_presentation_win32(slides, settings, output_file)
    else:
        return create_presentation_python_pptx(slides, settings, output_file)

def generate_readme(output_file="README.md"):
    """Generate a README file with syntax documentation.
    
    Args:
        output_file: Path to save the README file
        
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        readme_content = """# Markdown to PowerPoint Converter

This tool converts specially formatted Markdown files into feature-rich PowerPoint presentations. It supports a wide range of PowerPoint features including precise element positioning, formatting, animations, and more.

## Installation

### Requirements

- Python 3.6 or higher
- Required Python packages:
  - python-pptx
  - pyyaml

Install the dependencies with:

```bash
pip install python-pptx pyyaml
```

For advanced features (animations, transitions, etc.):
```bash
pip install pywin32
```

## Basic Usage

```bash
python md2pptx.py input.md -o output.pptx
```

Options:

- `-o, --output`: Specify the output PowerPoint file name
- `-b, --backend`: Choose backend (`win32com` or `python-pptx`, default: `win32com`)
- `-g, --generate-readme`: Generate a README file with syntax documentation

## Markdown Syntax Reference

The system uses a specialized Markdown format with YAML-style properties and custom element tags to describe PowerPoint features.

### Global Settings

At the top of your Markdown file, you can specify global presentation settings using YAML frontmatter:

```markdown
---
slide_width: 10in
slide_height: 7.5in
default_font: Arial
default_font_size: 18
default_font_color: "#333333"
company_logo: "path/to/logo.png"
footer_text: "Confidential"
---
```

### Slide Structure

Each slide begins with a level 1 or 2 heading:

```markdown
# Slide Title
```

### Slide Settings

You can specify slide-specific settings immediately after the slide title:

```markdown
# Slide Title

{
layout: "Title and Content",
background: "#FFFFFF",
transition: "fade",
footer: true,
slide_number: true
}
```

### Elements

Elements are defined using triple colon syntax followed by the element type:

```markdown
:::text
This is a text box
:::
```

Elements can have properties specified in brackets or in a separate property block:

```markdown
:::text[x:1in, y:2in, width:4in, height:2in]
This is a positioned text box
:::

# OR

:::text
{x:1in, y:2in, width:4in, height:2in, font_size:24, font_color:blue}
This is a styled text box
:::
```

## Element Types

### Text

Text elements create text boxes that can be precisely positioned and formatted.

```markdown
:::text
{
x: 1in,
y: 1in,
width: 8in,
height: 2in,
font: "Arial",
font_size: 24,
font_color: "#FF0000",
align: center,
bold: true,
italic: false,
underline: false
}
This is a styled text box with formatting applied.
Multiple paragraphs are supported.
:::
```

### Images

Image elements insert pictures from local or network files.

```markdown
:::image
{
x: 2in,
y: 3in,
width: 4in,
height: 3in,
rotation: 15,
transparency: 0.2,
border_color: "#000000",
border_width: 2pt
}
/path/to/image.jpg
:::
```

### Shapes

Shape elements create geometric shapes that can contain text.

```markdown
:::shape
{
x: 1in,
y: 1in,
width: 2in,
height: 2in,
shape_type: oval,
fill: "#0000FF",
border_color: "#000000",
border_width: 2pt,
transparency: 0.3,
rotation: 45,
shadow: true
}
Text inside shape
:::
```

Available shape types:

- rectangle
- oval
- rounded_rectangle
- triangle
- right_triangle
- diamond
- pentagon
- hexagon
- star
- arrow

### Tables

Table elements create data tables with rows and columns.

```markdown
:::table
{
x: 1in,
y: 1in,
width: 8in,
height: 3in,
header_row: true,
first_column: true,
banded_rows: true,
border_color: "#333333"
}
Header 1, Header 2, Header 3
Cell 1, Cell 2, Cell 3
Cell 4, Cell 5, Cell 6
:::
```

Or using Markdown table syntax:

```markdown
:::table
{
x: 1in,
y: 1in,
width: 8in,
height: 3in
}
| Header 1 | Header 2 | Header 3 |
| Cell 1 | Cell 2 | Cell 3 |
| Cell 4 | Cell 5 | Cell 6 |
:::
```

### Charts

Chart elements create data visualizations.

```markdown
:::chart
{
x: 1in,
y: 1in,
width: 6in,
height: 4in,
chart_type: bar,
title: "Sales Data",
has_legend: true,
legend_position: "bottom",
x_axis_title: "Categories",
y_axis_title: "Values"
}
Category 1, 10
Category 2, 15
Category 3, 7
Category 4, 12
:::
```

Available chart types:

- bar
- column
- line
- pie
- area
- scatter
- doughnut
- radar

### SmartArt

SmartArt elements create diagrams from lists of text items.

```markdown
:::smartart
{
x: 1in,
y: 1in,
width: 8in,
height: 3in,
type: "process"
}
Step 1, Step 2, Step 3, Step 4
:::
```

Available SmartArt types:

- process
- cycle
- hierarchy
- pyramid
- radial
- venn
- matrix
- relationship
- list

## Animation & Transition Effects

You can add animations to any element by adding animation properties:

```markdown
:::text
{
x: 1in,
y: 1in,
width: 4in,
animation: "fade",
animation_trigger: "on_click",
animation_delay: 0.5,
animation_duration: 1.0
}
This text will fade in when clicked.
:::
```

Slide transitions are specified in the slide settings:

```markdown
# Example Slide

{
transition: "fade",
transition_speed: "medium"
}
```

## Property Reference

### Global Settings Properties

| Property             | Description                        | Example Values                      |
| -------------------- | ---------------------------------- | ----------------------------------- |
| `slide_width`        | Width of all slides                | "10in", "25.4cm"                    |
| `slide_height`       | Height of all slides               | "7.5in", "19.05cm"                  |
| `default_font`       | Default font for all text          | "Arial", "Calibri"                  |
| `default_font_size`  | Default font size for all text     | 18, "18pt"                          |
| `default_font_color` | Default font color                 | "#333333", "rgb(51,51,51)", "black" |
| `company_logo`       | Path to company logo for templates | "/path/to/logo.png"                 |
| `footer_text`        | Default footer text for all slides | "Confidential"                      |
| `header_text`        | Default header text for all slides | "Company Name"                      |
| `template`           | PowerPoint template to use         | "template.potx"                     |

### Slide Properties

| Property           | Description                    | Example Values                              |
| ------------------ | ------------------------------ | ------------------------------------------- |
| `layout`           | Slide layout name              | "Title Slide", "Title and Content", "Blank" |
| `background`       | Background color or image path | "#FFFFFF", "/path/to/background.jpg"        |
| `background_style` | Style for background           | "solid", "gradient", "pattern"              |
| `transition`       | Slide transition effect        | "fade", "push", "wipe", "split"             |
| `transition_speed` | Speed of transition            | "slow", "medium", "fast"                    |
| `footer`           | Show/hide footer               | true, false                                 |
| `header`           | Show/hide header               | true, false                                 |
| `slide_number`     | Show/hide slide number         | true, false                                 |
| `notes`            | Speaker notes for the slide    | "Remember to mention key points"            |

### Animation Properties

| Property              | Description                      | Example Values                                     |
| --------------------- | -------------------------------- | -------------------------------------------------- |
| `animation`           | Animation type                   | "fade", "wipe", "fly_in", "float", "split", "zoom" |
| `animation_trigger`   | When animation starts            | "on_click", "with_previous", "after_previous"      |
| `animation_direction` | Direction of animation           | "in", "out", "up", "down", "left", "right"         |
| `animation_delay`     | Delay before animation (seconds) | 0, 0.5, 1                                          |
| `animation_duration`  | Animation duration (seconds)     | 0.5, 1, 2                                          |

## Element Properties

### Position and Size Properties

| Property   | Description         | Example Values            |
| ---------- | ------------------- | ------------------------- |
| `x`        | Horizontal position | "1in", "2.54cm", "72pt"   |
| `y`        | Vertical position   | "1in", "2.54cm", "72pt"   |
| `width`    | Element width       | "4in", "10.16cm", "288pt" |
| `height`   | Element height      | "3in", "7.62cm", "216pt"  |
| `rotation` | Rotation angle      | 45, 90, -30               |

### Working with Units

The system supports various units for position and size specifications:

- `in`: Inches (default if no unit specified)
- `cm`: Centimeters
- `pt`: Points (1/72 of an inch)
- `px`: Pixels (approximate conversion)

### Working with Colors

Colors can be specified in multiple formats:

- Hex: `#FF0000`, `#F00`
- Named colors: `red`, `blue`, `green`, etc.

## Tips and Best Practices

1. **Organizing Content**:
   - Use global settings for consistent presentation styles
   - Group related slides with meaningful titles
   - Use bullet points for lists and organized content

2. **Layout and Positioning**:
   - Use precise coordinates for complex layouts
   - Test your presentation with different screen resolutions
   - Consider using relative positioning for better adaptability

3. **Styling**:
   - Maintain consistent fonts and colors throughout
   - Use company style guides for visual elements
   - Limit animations to avoid distractions

4. **Performance**:
   - Optimize image sizes before embedding
   - Be mindful of the number of elements per slide
   - Split complex presentations into multiple files

## Troubleshooting

### Common Issues

1. **Element not appearing at expected position**:
   - Check units (in, cm, pt) on coordinates
   - Verify slide dimensions match expected size
   - Ensure element is not positioned outside slide boundaries

2. **Formatting not applied**:
   - Check syntax of property specifications
   - Verify property names are spelled correctly
   - Ensure values are in expected formats

3. **Images not displaying**:
   - Verify file paths are correct
   - Check that image files exist and are readable
   - Try using absolute paths if relative paths fail

## Advanced Features

These features are available when using the win32com backend:

- **Animation effects** for element entrance, emphasis, and exit
- **Slide transitions** like fade, push, wipe, etc.
- **SmartArt** diagrams for visual representation of relationships
- **Custom templates** for consistent branding
- **Speaker notes** for presentation guidance

## Contributing

Contributions to improve the conversion tool are welcome. Please submit issues and pull requests to the project repository.

## License

This tool is licensed under the MIT License.
"""

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(readme_content)
        
        print(f"Generated README at: {os.path.abspath(output_file)}")
        return True
    
    except Exception as e:
        print(f"Error generating README: {e}")
        return False

def main():
    """Main function to process command line arguments and execute the conversion."""
    parser = argparse.ArgumentParser(description='Convert Markdown to PowerPoint')
    parser.add_argument('input_file', nargs='?', help='Input Markdown file')
    parser.add_argument('-o', '--output', help='Output PowerPoint file')
    parser.add_argument('-b', '--backend', choices=['win32com', 'python-pptx'], 
                      default='win32com', help='Backend to use for PowerPoint generation')
    parser.add_argument('-g', '--generate-readme', action='store_true', 
                      help='Generate README with syntax documentation')
    
    args = parser.parse_args()
    
    # Handle README generation
    if args.generate_readme:
        generate_readme()
        if not args.input_file:
            return 0
    
    # Validate input file
    if not args.input_file:
        parser.print_help()
        return 1
        
    if not os.path.exists(args.input_file):
        print(f"Error: Input file '{args.input_file}' not found.")
        return 1
    
    # Determine output file
    output_file = args.output
    if not output_file:
        input_path = Path(args.input_file)
        output_file = str(input_path.with_suffix('.pptx'))
    
    # Parse markdown
    print(f"Parsing {args.input_file}...")
    slides, settings = parse_markdown(args.input_file)
    print(f"Found {len(slides)} slides")
    
    # Set backend
    backend = BackendType.WIN32COM if args.backend == 'win32com' else BackendType.PYTHON_PPTX
    
    # Create PowerPoint
    if create_presentation(slides, settings, output_file, backend):
        print("\nSUCCESS: Converted markdown to PowerPoint")
        return 0
    else:
        print("\nFAILED: Could not create PowerPoint presentation")
        return 1

if __name__ == "__main__":
    sys.exit(main())